from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.core.extraction_runtime import ExtractionRuntimeRequest
from app.core.local_extraction import (
    ERROR_CODE_LOCAL_DOCX_EMPTY,
    ERROR_CODE_LOCAL_PDF_TEXT_LAYER_EMPTY,
    ERROR_CODE_LOCAL_PPTX_EMPTY,
    ERROR_CODE_LOCAL_SOURCE_NOT_FOUND,
    ERROR_CODE_LOCAL_UNSUPPORTED_FILE_TYPE,
    ERROR_CODE_LOCAL_XLSX_EMPTY,
    LOCAL_DOCX_PROFILE_NAME,
    LOCAL_MARKDOWN_PROFILE_NAME,
    LOCAL_PDF_TEXT_PROFILE_NAME,
    LOCAL_PLAIN_TEXT_PROFILE_NAME,
    LOCAL_PPTX_PROFILE_NAME,
    LOCAL_XLSX_PROFILE_NAME,
    SUPPORTED_LOCAL_PROFILE_SUFFIXES,
    get_local_extraction_handler,
    list_local_extraction_handlers,
    normalize_file_type,
    run_local_extraction,
    select_local_extraction_handler,
    select_local_extraction_profile_name,
)


def make_minimal_pdf(lines: list[str] | None = None) -> bytes:
    text_operations: list[str] = []
    y_position = 760
    for line in lines or []:
        escaped = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        text_operations.append(f"BT /F1 12 Tf 72 {y_position} Td ({escaped}) Tj ET")
        y_position -= 18

    stream = "\n".join(text_operations).encode("ascii")
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        (
            b"3 0 obj\n"
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>\n"
            b"endobj\n"
        ),
        b"4 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
        (
            b"5 0 obj\n<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"\nendstream\nendobj\n"
        ),
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for pdf_object in objects:
        offsets.append(len(output))
        output.extend(pdf_object)

    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        (
            f"trailer\n<< /Root 1 0 R /Size {len(objects) + 1} >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(output)


def write_sample_docx(path: Path) -> None:
    document = Document()
    document.add_heading("DOCX Title", level=1)
    document.add_paragraph("First paragraph from DOCX.")
    document.add_heading("Data", level=2)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Accuracy"
    table.cell(1, 1).text = "High"
    document.save(path)


def write_sample_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PPTX Title"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    textbox.text_frame.text = "First paragraph from PPTX."

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(6), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Quality"
    table.cell(1, 1).text = "Baseline"
    presentation.save(path)


def write_sample_xlsx(path: Path) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Measurements"
    worksheet.append(["Metric", "Value"])
    worksheet.append(["Accuracy", "High"])
    worksheet.append(["Latency", "Low"])
    workbook.save(path)


def test_local_extraction_registry_selects_implemented_profiles() -> None:
    handlers = {handler.profile_name: handler for handler in list_local_extraction_handlers()}

    markdown_handler = get_local_extraction_handler(LOCAL_MARKDOWN_PROFILE_NAME)
    pdf_handler = get_local_extraction_handler(LOCAL_PDF_TEXT_PROFILE_NAME)
    docx_handler = get_local_extraction_handler(LOCAL_DOCX_PROFILE_NAME)
    pptx_handler = get_local_extraction_handler(LOCAL_PPTX_PROFILE_NAME)
    xlsx_handler = get_local_extraction_handler(LOCAL_XLSX_PROFILE_NAME)

    assert set(handlers) == {
        LOCAL_MARKDOWN_PROFILE_NAME,
        LOCAL_PLAIN_TEXT_PROFILE_NAME,
        LOCAL_PDF_TEXT_PROFILE_NAME,
        LOCAL_DOCX_PROFILE_NAME,
        LOCAL_PPTX_PROFILE_NAME,
        LOCAL_XLSX_PROFILE_NAME,
    }
    assert markdown_handler is not None
    assert pdf_handler is not None
    assert docx_handler is not None
    assert pptx_handler is not None
    assert xlsx_handler is not None
    assert markdown_handler.supports_file_type(".MD")
    assert pdf_handler.supports_file_type(".PDF")
    assert docx_handler.supports_file_type(".DOCX")
    assert pptx_handler.supports_file_type(".PPTX")
    assert xlsx_handler.supports_file_type(".XLSX")
    assert normalize_file_type(".Txt") == "txt"
    assert select_local_extraction_handler("md") == markdown_handler
    assert select_local_extraction_profile_name(".text") == LOCAL_PLAIN_TEXT_PROFILE_NAME
    assert select_local_extraction_profile_name("pdf") == LOCAL_PDF_TEXT_PROFILE_NAME
    assert select_local_extraction_profile_name("docx") == LOCAL_DOCX_PROFILE_NAME
    assert select_local_extraction_profile_name("pptx") == LOCAL_PPTX_PROFILE_NAME
    assert select_local_extraction_profile_name("xlsx") == LOCAL_XLSX_PROFILE_NAME
    assert select_local_extraction_handler("pdf") == pdf_handler
    assert select_local_extraction_handler("pptx") == pptx_handler
    assert select_local_extraction_handler("xlsx") == xlsx_handler
    assert get_local_extraction_handler("missing_profile") is None
    assert SUPPORTED_LOCAL_PROFILE_SUFFIXES[LOCAL_MARKDOWN_PROFILE_NAME] == {".md"}
    assert SUPPORTED_LOCAL_PROFILE_SUFFIXES[LOCAL_PDF_TEXT_PROFILE_NAME] == {".pdf"}
    assert SUPPORTED_LOCAL_PROFILE_SUFFIXES[LOCAL_DOCX_PROFILE_NAME] == {".docx"}
    assert SUPPORTED_LOCAL_PROFILE_SUFFIXES[LOCAL_PPTX_PROFILE_NAME] == {".pptx"}
    assert SUPPORTED_LOCAL_PROFILE_SUFFIXES[LOCAL_XLSX_PROFILE_NAME] == {".xlsx"}


def test_run_local_markdown_extraction_returns_artifact_and_blocks(tmp_path: Path) -> None:
    source_path = tmp_path / "sample.md"
    source_path.write_text(
        "# Title\r\n\r\n"
        "Intro paragraph.\r\n\r\n"
        "## Data\r\n\r\n"
        "| A | B |\r\n"
        "| --- | --- |\r\n"
        "| 1 | 2 |\r\n\r\n"
        "```python\r\n"
        "print(1)\r\n"
        "```\r\n",
        encoding="utf-8",
    )

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_MARKDOWN_PROFILE_NAME,
            mime_type="text/markdown",
            detected_file_type="md",
        )
    )

    assert result.status == "succeeded"
    assert result.errors == ()
    assert result.artifacts[0].artifact_type == "normalized_markdown"
    assert "\r" not in result.artifacts[0].content_text
    assert result.artifacts[0].metadata["parser_name"] == "markdown"
    assert [block.block_type for block in result.blocks] == [
        "heading",
        "paragraph",
        "heading",
        "table",
        "code",
    ]
    assert result.blocks[1].parent_block_seq == 0
    assert result.blocks[2].parent_block_seq == 0
    assert result.blocks[3].parent_block_seq == 2
    assert result.blocks[3].source_anchor == {
        "start_line": 7,
        "end_line": 9,
        "table_index": 3,
    }
    assert result.blocks[4].metadata == {"language": "python"}
    assert result.runtime_metadata["block_count"] == 5


def test_run_local_extraction_uses_detected_file_type_when_suffix_is_missing(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "stored-file"
    source_path.write_text("# Title\n\nBody.", encoding="utf-8")

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_MARKDOWN_PROFILE_NAME,
            detected_file_type="md",
        )
    )

    assert result.status == "succeeded"
    assert result.artifacts[0].artifact_type == "normalized_markdown"
    assert result.blocks[0].block_type == "heading"


def test_run_local_plain_text_extraction_splits_paragraphs(tmp_path: Path) -> None:
    source_path = tmp_path / "sample.txt"
    source_path.write_text(
        "First paragraph\r\ncontinues.\r\n\r\nSecond paragraph.",
        encoding="utf-8",
    )

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_PLAIN_TEXT_PROFILE_NAME,
            detected_file_type="txt",
        )
    )

    assert result.status == "succeeded"
    assert result.artifacts[0].artifact_type == "plain_text"
    assert result.blocks[0].content_text == "First paragraph continues."
    assert result.blocks[0].source_anchor == {"start_line": 1, "end_line": 2}
    assert result.blocks[1].content_text == "Second paragraph."
    assert result.runtime_metadata["extractor_name"] == "local_plain_text"


def test_run_local_pdf_text_extraction_returns_artifact_and_page_blocks(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "sample.pdf"
    source_path.write_bytes(
        make_minimal_pdf(
            [
                "PDF Title",
                "",
                "First paragraph from the text layer.",
                "Second paragraph continues.",
            ]
        )
    )

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_PDF_TEXT_PROFILE_NAME,
            mime_type="application/pdf",
            detected_file_type="pdf",
        )
    )

    assert result.status == "succeeded"
    assert result.artifacts[0].artifact_type == "normalized_markdown"
    assert "<!-- page: 1 -->" in result.artifacts[0].content_text
    assert result.artifacts[0].metadata["library"] == "pypdf"
    assert result.artifacts[0].metadata["page_count"] == 1
    assert result.artifacts[0].metadata["text_layer_only"] is True
    assert result.artifacts[0].metadata["ocr_enabled"] is False
    assert [block.page_no for block in result.blocks] == [1]
    assert result.blocks[0].content_text.startswith("PDF Title")
    assert result.blocks[0].heading_path == ("Page 1",)
    assert result.blocks[0].source_anchor == {"page_no": 1, "paragraph_index": 1}
    assert result.blocks[0].metadata == {
        "source": "pdf_text_layer",
        "library": "pypdf",
    }
    assert result.runtime_metadata["extractor_name"] == "local_pdf_text"
    assert result.runtime_metadata["page_count"] == 1


def test_run_local_pdf_text_extraction_fails_without_text_layer(tmp_path: Path) -> None:
    source_path = tmp_path / "empty.pdf"
    source_path.write_bytes(make_minimal_pdf())

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_PDF_TEXT_PROFILE_NAME,
            mime_type="application/pdf",
            detected_file_type="pdf",
        )
    )

    assert result.status == "failed"
    assert result.runtime_metadata["error_code"] == ERROR_CODE_LOCAL_PDF_TEXT_LAYER_EMPTY
    assert "text layer" in result.errors[0]


def test_run_local_docx_extraction_returns_artifact_headings_and_table_blocks(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "sample.docx"
    write_sample_docx(source_path)

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_DOCX_PROFILE_NAME,
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            detected_file_type="docx",
        )
    )

    assert result.status == "succeeded"
    assert result.artifacts[0].artifact_type == "normalized_markdown"
    assert "# DOCX Title" in result.artifacts[0].content_text
    assert "| Metric | Value |" in result.artifacts[0].content_text
    assert result.artifacts[0].metadata["library"] == "python-docx"
    assert result.artifacts[0].metadata["paragraph_count"] == 1
    assert result.artifacts[0].metadata["heading_count"] == 2
    assert result.artifacts[0].metadata["table_count"] == 1
    assert [block.block_type for block in result.blocks] == [
        "heading",
        "paragraph",
        "heading",
        "table",
    ]
    assert result.blocks[0].content_markdown == "# DOCX Title"
    assert result.blocks[1].parent_block_seq == 0
    assert result.blocks[1].heading_path == ("DOCX Title",)
    assert result.blocks[2].parent_block_seq == 0
    assert result.blocks[2].heading_path == ("DOCX Title", "Data")
    assert result.blocks[3].parent_block_seq == 2
    assert result.blocks[3].heading_path == ("DOCX Title", "Data")
    assert result.blocks[3].source_anchor == {
        "body_index": 3,
        "table_index": 1,
    }
    assert result.blocks[3].metadata == {
        "source": "docx",
        "row_count": 2,
        "column_count": 2,
    }
    assert result.runtime_metadata["extractor_name"] == "local_docx"
    assert result.runtime_metadata["table_count"] == 1


def test_run_local_docx_extraction_fails_without_extractable_text(tmp_path: Path) -> None:
    source_path = tmp_path / "empty.docx"
    Document().save(source_path)

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_DOCX_PROFILE_NAME,
            detected_file_type="docx",
        )
    )

    assert result.status == "failed"
    assert result.runtime_metadata["error_code"] == ERROR_CODE_LOCAL_DOCX_EMPTY
    assert "DOCX does not contain" in result.errors[0]


def test_run_local_pptx_extraction_returns_artifact_slide_text_and_table_blocks(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "sample.pptx"
    write_sample_pptx(source_path)

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_PPTX_PROFILE_NAME,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            detected_file_type="pptx",
        )
    )

    assert result.status == "succeeded"
    assert result.artifacts[0].artifact_type == "normalized_markdown"
    assert "<!-- slide: 1 -->" in result.artifacts[0].content_text
    assert "# PPTX Title" in result.artifacts[0].content_text
    assert "| Metric | Value |" in result.artifacts[0].content_text
    assert result.artifacts[0].metadata["library"] == "python-pptx"
    assert result.artifacts[0].metadata["slide_count"] == 1
    assert result.artifacts[0].metadata["text_shape_count"] == 2
    assert result.artifacts[0].metadata["table_count"] == 1
    assert result.artifacts[0].metadata["preserve_slide_boundaries"] is True
    assert [block.block_type for block in result.blocks] == [
        "heading",
        "paragraph",
        "table",
    ]
    assert [block.slide_no for block in result.blocks] == [1, 1, 1]
    assert result.blocks[0].content_markdown == "# PPTX Title"
    assert result.blocks[0].heading_path == ("PPTX Title",)
    assert result.blocks[0].source_anchor["slide_no"] == 1
    assert result.blocks[1].parent_block_seq == 0
    assert result.blocks[1].heading_path == ("PPTX Title",)
    assert result.blocks[2].parent_block_seq == 0
    assert result.blocks[2].heading_path == ("PPTX Title",)
    assert result.blocks[2].metadata["source"] == "pptx"
    assert result.blocks[2].metadata["row_count"] == 2
    assert result.blocks[2].metadata["column_count"] == 2
    assert result.blocks[2].metadata["shape_name"].startswith("Table")
    assert result.blocks[2].source_anchor["table_index"] == 1
    assert result.runtime_metadata["extractor_name"] == "local_pptx"
    assert result.runtime_metadata["table_count"] == 1


def test_run_local_pptx_extraction_fails_without_extractable_text(tmp_path: Path) -> None:
    source_path = tmp_path / "empty.pptx"
    Presentation().save(source_path)

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_PPTX_PROFILE_NAME,
            detected_file_type="pptx",
        )
    )

    assert result.status == "failed"
    assert result.runtime_metadata["error_code"] == ERROR_CODE_LOCAL_PPTX_EMPTY
    assert "PPTX does not contain" in result.errors[0]


def test_run_local_xlsx_extraction_returns_artifact_sheet_heading_and_table_blocks(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "sample.xlsx"
    write_sample_xlsx(source_path)

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_XLSX_PROFILE_NAME,
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            detected_file_type="xlsx",
        )
    )

    assert result.status == "succeeded"
    assert result.artifacts[0].artifact_type == "normalized_markdown"
    assert "<!-- sheet: Measurements -->" in result.artifacts[0].content_text
    assert "# Measurements" in result.artifacts[0].content_text
    assert "| Metric | Value |" in result.artifacts[0].content_text
    assert result.artifacts[0].metadata["library"] == "openpyxl"
    assert result.artifacts[0].metadata["sheet_count"] == 1
    assert result.artifacts[0].metadata["extracted_sheet_count"] == 1
    assert result.artifacts[0].metadata["table_count"] == 1
    assert result.artifacts[0].metadata["cell_count"] == 6
    assert result.artifacts[0].metadata["emit_markdown_tables"] is True
    assert [block.block_type for block in result.blocks] == ["heading", "table"]
    assert [block.sheet_name for block in result.blocks] == ["Measurements", "Measurements"]
    assert result.blocks[0].content_markdown == "# Measurements"
    assert result.blocks[0].heading_path == ("Measurements",)
    assert result.blocks[1].parent_block_seq == 0
    assert result.blocks[1].heading_path == ("Measurements",)
    assert result.blocks[1].cell_range == "A1:B3"
    assert result.blocks[1].source_anchor == {
        "sheet_index": 1,
        "sheet_name": "Measurements",
        "cell_range": "A1:B3",
        "table_index": 1,
    }
    assert result.blocks[1].metadata == {
        "source": "xlsx",
        "row_count": 3,
        "column_count": 2,
        "cell_count": 6,
    }
    assert result.runtime_metadata["extractor_name"] == "local_xlsx"
    assert result.runtime_metadata["table_count"] == 1


def test_run_local_xlsx_extraction_fails_without_extractable_cells(tmp_path: Path) -> None:
    source_path = tmp_path / "empty.xlsx"
    Workbook().save(source_path)

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            document_id=2,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_XLSX_PROFILE_NAME,
            detected_file_type="xlsx",
        )
    )

    assert result.status == "failed"
    assert result.runtime_metadata["error_code"] == ERROR_CODE_LOCAL_XLSX_EMPTY
    assert "XLSX does not contain" in result.errors[0]


def test_run_local_extraction_returns_failed_for_missing_source(tmp_path: Path) -> None:
    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            storage_path=str(tmp_path / "missing.md"),
            extraction_profile_name=LOCAL_MARKDOWN_PROFILE_NAME,
        )
    )

    assert result.status == "failed"
    assert result.runtime_metadata["error_code"] == ERROR_CODE_LOCAL_SOURCE_NOT_FOUND
    assert "not found" in result.errors[0]


def test_run_local_extraction_returns_failed_for_profile_suffix_mismatch(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "sample.pdf"
    source_path.write_text("# Not really PDF", encoding="utf-8")

    result = run_local_extraction(
        ExtractionRuntimeRequest(
            file_id=1,
            storage_path=str(source_path),
            extraction_profile_name=LOCAL_MARKDOWN_PROFILE_NAME,
        )
    )

    assert result.status == "failed"
    assert result.runtime_metadata["error_code"] == ERROR_CODE_LOCAL_UNSUPPORTED_FILE_TYPE
    assert "cannot process" in result.errors[0]
