"""Deterministic local extraction fixture corpus and stable snapshots."""

from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.core.extraction_runtime import ExtractionRuntimeResult
from app.core.local_extraction import (
    LOCAL_DOCX_PROFILE_NAME,
    LOCAL_HWPX_PROFILE_NAME,
    LOCAL_MARKDOWN_PROFILE_NAME,
    LOCAL_PDF_TEXT_PROFILE_NAME,
    LOCAL_PLAIN_TEXT_PROFILE_NAME,
    LOCAL_PPTX_PROFILE_NAME,
    LOCAL_XLSX_PROFILE_NAME,
)


@dataclass(frozen=True)
class ExtractionFixtureCase:
    case_id: str
    file_name: str
    profile_name: str
    detected_file_type: str
    mime_type: str | None
    write_source: Callable[[Path], None]
    expected_snapshot: dict[str, Any]


STABLE_METADATA_EXCLUDED_KEYS = {
    "library_version",
    "options",
    "shape_name",
    "source_file_name",
    "source_path",
}

MARKDOWN_FIXTURE_TEXT = (
    "# Fixture Markdown\n\n"
    "Intro paragraph.\n\n"
    "## Metrics\n\n"
    "| Metric | Value |\n"
    "| --- | --- |\n"
    "| Accuracy | High |\n"
)

PLAIN_TEXT_FIXTURE_TEXT = "First plain paragraph.\n" "continues.\n\n" "Second plain paragraph.\n"


def write_markdown_fixture(path: Path) -> None:
    path.write_text(MARKDOWN_FIXTURE_TEXT, encoding="utf-8")


def write_plain_text_fixture(path: Path) -> None:
    path.write_text(PLAIN_TEXT_FIXTURE_TEXT, encoding="utf-8")


def write_pdf_fixture(path: Path) -> None:
    path.write_bytes(make_minimal_pdf(["PDF Fixture paragraph."]))


def write_docx_fixture(path: Path) -> None:
    document = Document()
    document.add_heading("DOCX Fixture", level=1)
    document.add_paragraph("DOCX body paragraph.")
    document.add_heading("Metrics", level=2)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Accuracy"
    table.cell(1, 1).text = "High"
    document.save(path)


def write_pptx_fixture(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PPTX Fixture"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    textbox.text_frame.text = "PPTX body paragraph."

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(6), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Quality"
    table.cell(1, 1).text = "Baseline"
    presentation.save(path)


def write_xlsx_fixture(path: Path) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Fixture Sheet"
    worksheet.append(["Metric", "Value"])
    worksheet.append(["Accuracy", "High"])
    worksheet.append(["Latency", "Low"])
    workbook.save(path)


def write_hwpx_fixture(path: Path) -> None:
    section_xml = """<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p>
    <hp:run><hp:t>HWPX body paragraph.</hp:t></hp:run>
  </hp:p>
  <hp:tbl>
    <hp:tr>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Metric</hp:t></hp:run></hp:p></hp:subList></hp:tc>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Value</hp:t></hp:run></hp:p></hp:subList></hp:tc>
    </hp:tr>
    <hp:tr>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Quality</hp:t></hp:run></hp:p></hp:subList></hp:tc>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Baseline</hp:t></hp:run></hp:p></hp:subList></hp:tc>
    </hp:tr>
  </hp:tbl>
</hp:sec>
"""
    with ZipFile(path, "w") as archive:
        archive.writestr("mimetype", "application/hwp+zip")
        archive.writestr("Contents/section0.xml", section_xml)


def make_minimal_pdf(lines: list[str] | None = None) -> bytes:
    text_operations: list[str] = []
    y_position = 760
    for line in lines or []:
        escaped = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        text_operations.append(f"BT /F1 12 Tf 72 {y_position} Td ({escaped}) Tj ET")
        y_position -= 18

    stream = "\n".join(text_operations).encode("ascii")
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        (
            b"3 0 obj\n"
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>\n"
            b"endobj\n"
        ),
        b"4 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
        (
            b"5 0 obj\n<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"\nendstream\nendobj\n"
        ),
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for pdf_object in objects:
        offsets.append(len(output))
        output.extend(pdf_object)

    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        (
            f"trailer\n<< /Root 1 0 R /Size {len(objects) + 1} >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(output)


def extraction_result_snapshot(result: ExtractionRuntimeResult) -> dict[str, Any]:
    artifact = result.artifacts[0] if result.artifacts else None
    payload: dict[str, Any] = {
        "status": result.status,
        "runtime_metadata": stable_payload(result.runtime_metadata),
        "warnings": list(result.warnings),
        "errors": list(result.errors),
        "blocks": [block_snapshot(block) for block in result.blocks],
    }
    if artifact is not None:
        payload["artifact"] = compact_payload(
            {
                "artifact_type": artifact.artifact_type,
                "content_text": artifact.content_text,
                "language": artifact.language,
                "metadata": stable_payload(artifact.metadata),
            }
        )
    return compact_payload(payload)


def block_snapshot(block: Any) -> dict[str, Any]:
    return compact_payload(
        {
            "block_seq": block.block_seq,
            "block_type": block.block_type,
            "parent_block_seq": block.parent_block_seq,
            "content_text": block.content_text,
            "content_markdown": block.content_markdown,
            "heading_path": list(block.heading_path),
            "source_anchor": stable_payload(block.source_anchor),
            "page_no": block.page_no,
            "slide_no": block.slide_no,
            "sheet_name": block.sheet_name,
            "cell_range": block.cell_range,
            "metadata": stable_payload(block.metadata),
        }
    )


def stable_payload(value: Any) -> Any:
    if isinstance(value, dict):
        return compact_payload(
            {
                key: stable_payload(item)
                for key, item in value.items()
                if key not in STABLE_METADATA_EXCLUDED_KEYS
            }
        )
    if isinstance(value, tuple | list):
        return [stable_payload(item) for item in value]
    return value


def compact_payload(payload: dict[str, Any]) -> dict[str, Any]:
    return {
        key: value
        for key, value in payload.items()
        if value is not None and value != {} and value != [] and value != ()
    }


EXTRACTION_FIXTURE_CASES = (
    ExtractionFixtureCase(
        case_id="markdown",
        file_name="fixture.md",
        profile_name=LOCAL_MARKDOWN_PROFILE_NAME,
        detected_file_type="md",
        mime_type="text/markdown",
        write_source=write_markdown_fixture,
        expected_snapshot={
            "status": "succeeded",
            "runtime_metadata": {
                "profile_name": LOCAL_MARKDOWN_PROFILE_NAME,
                "extractor_name": "markdown",
                "extractor_version": "0.1.0",
                "line_count": 9,
                "block_count": 4,
            },
            "artifact": {
                "artifact_type": "normalized_markdown",
                "content_text": MARKDOWN_FIXTURE_TEXT,
                "metadata": {
                    "parser_name": "markdown",
                    "parser_version": "0.1.0",
                    "line_count": 9,
                    "block_count": 4,
                },
            },
            "blocks": [
                {
                    "block_seq": 0,
                    "block_type": "heading",
                    "content_text": "Fixture Markdown",
                    "content_markdown": "# Fixture Markdown",
                    "heading_path": ["Fixture Markdown"],
                    "source_anchor": {"start_line": 1, "end_line": 1},
                    "metadata": {"level": 1},
                },
                {
                    "block_seq": 1,
                    "block_type": "paragraph",
                    "parent_block_seq": 0,
                    "content_text": "Intro paragraph.",
                    "content_markdown": "Intro paragraph.",
                    "heading_path": ["Fixture Markdown"],
                    "source_anchor": {"start_line": 3, "end_line": 3},
                },
                {
                    "block_seq": 2,
                    "block_type": "heading",
                    "parent_block_seq": 0,
                    "content_text": "Metrics",
                    "content_markdown": "## Metrics",
                    "heading_path": ["Fixture Markdown", "Metrics"],
                    "source_anchor": {"start_line": 5, "end_line": 5},
                    "metadata": {"level": 2},
                },
                {
                    "block_seq": 3,
                    "block_type": "table",
                    "parent_block_seq": 2,
                    "content_text": (
                        "| Metric | Value |\n" "| --- | --- |\n" "| Accuracy | High |"
                    ),
                    "content_markdown": (
                        "| Metric | Value |\n" "| --- | --- |\n" "| Accuracy | High |"
                    ),
                    "heading_path": ["Fixture Markdown", "Metrics"],
                    "source_anchor": {
                        "start_line": 7,
                        "end_line": 9,
                        "table_index": 3,
                    },
                    "metadata": {
                        "columns": ["Metric", "Value"],
                        "row_count": 1,
                    },
                },
            ],
        },
    ),
    ExtractionFixtureCase(
        case_id="plain_text",
        file_name="fixture.txt",
        profile_name=LOCAL_PLAIN_TEXT_PROFILE_NAME,
        detected_file_type="txt",
        mime_type="text/plain",
        write_source=write_plain_text_fixture,
        expected_snapshot={
            "status": "succeeded",
            "runtime_metadata": {
                "profile_name": LOCAL_PLAIN_TEXT_PROFILE_NAME,
                "extractor_name": "local_plain_text",
                "extractor_version": "0.1.0",
                "line_count": 4,
                "block_count": 2,
            },
            "artifact": {
                "artifact_type": "plain_text",
                "content_text": PLAIN_TEXT_FIXTURE_TEXT,
                "metadata": {"line_count": 4, "block_count": 2},
            },
            "blocks": [
                {
                    "block_seq": 0,
                    "block_type": "paragraph",
                    "content_text": "First plain paragraph. continues.",
                    "content_markdown": "First plain paragraph. continues.",
                    "source_anchor": {"start_line": 1, "end_line": 2},
                },
                {
                    "block_seq": 1,
                    "block_type": "paragraph",
                    "content_text": "Second plain paragraph.",
                    "content_markdown": "Second plain paragraph.",
                    "source_anchor": {"start_line": 4, "end_line": 4},
                },
            ],
        },
    ),
    ExtractionFixtureCase(
        case_id="pdf",
        file_name="fixture.pdf",
        profile_name=LOCAL_PDF_TEXT_PROFILE_NAME,
        detected_file_type="pdf",
        mime_type="application/pdf",
        write_source=write_pdf_fixture,
        expected_snapshot={
            "status": "succeeded",
            "runtime_metadata": {
                "profile_name": LOCAL_PDF_TEXT_PROFILE_NAME,
                "extractor_name": "local_pdf_text",
                "extractor_version": "0.1.0",
                "library": "pypdf",
                "page_count": 1,
                "extracted_page_count": 1,
                "block_count": 1,
                "text_layer_only": True,
                "ocr_enabled": False,
            },
            "artifact": {
                "artifact_type": "normalized_markdown",
                "content_text": "<!-- page: 1 -->\n\nPDF Fixture paragraph.",
                "metadata": {
                    "parser_name": "local_pdf_text",
                    "parser_version": "0.1.0",
                    "library": "pypdf",
                    "page_count": 1,
                    "extracted_page_count": 1,
                    "block_count": 1,
                    "text_layer_only": True,
                    "ocr_enabled": False,
                },
            },
            "blocks": [
                {
                    "block_seq": 0,
                    "block_type": "paragraph",
                    "content_text": "PDF Fixture paragraph.",
                    "content_markdown": "PDF Fixture paragraph.",
                    "heading_path": ["Page 1"],
                    "source_anchor": {"page_no": 1, "paragraph_index": 1},
                    "page_no": 1,
                    "metadata": {"source": "pdf_text_layer", "library": "pypdf"},
                }
            ],
        },
    ),
    ExtractionFixtureCase(
        case_id="docx",
        file_name="fixture.docx",
        profile_name=LOCAL_DOCX_PROFILE_NAME,
        detected_file_type="docx",
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        write_source=write_docx_fixture,
        expected_snapshot={
            "status": "succeeded",
            "runtime_metadata": {
                "profile_name": LOCAL_DOCX_PROFILE_NAME,
                "extractor_name": "local_docx",
                "extractor_version": "0.1.0",
                "library": "python-docx",
                "paragraph_count": 1,
                "heading_count": 2,
                "table_count": 1,
                "block_count": 4,
                "preserve_headings": True,
                "preserve_tables": True,
            },
            "artifact": {
                "artifact_type": "normalized_markdown",
                "content_text": (
                    "# DOCX Fixture\n\n"
                    "DOCX body paragraph.\n\n"
                    "## Metrics\n\n"
                    "| Metric | Value |\n"
                    "| --- | --- |\n"
                    "| Accuracy | High |"
                ),
                "metadata": {
                    "parser_name": "local_docx",
                    "parser_version": "0.1.0",
                    "library": "python-docx",
                    "paragraph_count": 1,
                    "heading_count": 2,
                    "table_count": 1,
                    "block_count": 4,
                    "preserve_headings": True,
                    "preserve_tables": True,
                },
            },
            "blocks": [
                {
                    "block_seq": 0,
                    "block_type": "heading",
                    "content_text": "DOCX Fixture",
                    "content_markdown": "# DOCX Fixture",
                    "heading_path": ["DOCX Fixture"],
                    "source_anchor": {
                        "body_index": 0,
                        "paragraph_index": 1,
                        "style_name": "Heading 1",
                    },
                    "metadata": {
                        "source": "docx",
                        "style_name": "Heading 1",
                        "heading_level": 1,
                    },
                },
                {
                    "block_seq": 1,
                    "block_type": "paragraph",
                    "parent_block_seq": 0,
                    "content_text": "DOCX body paragraph.",
                    "content_markdown": "DOCX body paragraph.",
                    "heading_path": ["DOCX Fixture"],
                    "source_anchor": {
                        "body_index": 1,
                        "paragraph_index": 2,
                        "style_name": "Normal",
                    },
                    "metadata": {
                        "source": "docx",
                        "style_name": "Normal",
                    },
                },
                {
                    "block_seq": 2,
                    "block_type": "heading",
                    "parent_block_seq": 0,
                    "content_text": "Metrics",
                    "content_markdown": "## Metrics",
                    "heading_path": ["DOCX Fixture", "Metrics"],
                    "source_anchor": {
                        "body_index": 2,
                        "paragraph_index": 3,
                        "style_name": "Heading 2",
                    },
                    "metadata": {
                        "source": "docx",
                        "style_name": "Heading 2",
                        "heading_level": 2,
                    },
                },
                {
                    "block_seq": 3,
                    "block_type": "table",
                    "parent_block_seq": 2,
                    "content_text": "Metric\tValue\nAccuracy\tHigh",
                    "content_markdown": (
                        "| Metric | Value |\n" "| --- | --- |\n" "| Accuracy | High |"
                    ),
                    "heading_path": ["DOCX Fixture", "Metrics"],
                    "source_anchor": {"body_index": 3, "table_index": 1},
                    "metadata": {
                        "source": "docx",
                        "row_count": 2,
                        "column_count": 2,
                    },
                },
            ],
        },
    ),
    ExtractionFixtureCase(
        case_id="pptx",
        file_name="fixture.pptx",
        profile_name=LOCAL_PPTX_PROFILE_NAME,
        detected_file_type="pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        write_source=write_pptx_fixture,
        expected_snapshot={
            "status": "succeeded",
            "runtime_metadata": {
                "profile_name": LOCAL_PPTX_PROFILE_NAME,
                "extractor_name": "local_pptx",
                "extractor_version": "0.1.0",
                "library": "python-pptx",
                "slide_count": 1,
                "text_shape_count": 2,
                "table_count": 1,
                "block_count": 3,
                "preserve_slide_boundaries": True,
                "preserve_tables": True,
            },
            "artifact": {
                "artifact_type": "normalized_markdown",
                "content_text": (
                    "<!-- slide: 1 -->\n\n"
                    "# PPTX Fixture\n\n"
                    "PPTX body paragraph.\n\n"
                    "| Metric | Value |\n"
                    "| --- | --- |\n"
                    "| Quality | Baseline |"
                ),
                "metadata": {
                    "parser_name": "local_pptx",
                    "parser_version": "0.1.0",
                    "library": "python-pptx",
                    "slide_count": 1,
                    "text_shape_count": 2,
                    "table_count": 1,
                    "block_count": 3,
                    "preserve_slide_boundaries": True,
                    "preserve_tables": True,
                },
            },
            "blocks": [
                {
                    "block_seq": 0,
                    "block_type": "heading",
                    "content_text": "PPTX Fixture",
                    "content_markdown": "# PPTX Fixture",
                    "heading_path": ["PPTX Fixture"],
                    "source_anchor": {
                        "slide_no": 1,
                        "shape_index": 1,
                        "text_shape_index": 1,
                        "slide_text_shape_index": 1,
                    },
                    "slide_no": 1,
                    "metadata": {"source": "pptx", "is_title": True},
                },
                {
                    "block_seq": 1,
                    "block_type": "paragraph",
                    "parent_block_seq": 0,
                    "content_text": "PPTX body paragraph.",
                    "content_markdown": "PPTX body paragraph.",
                    "heading_path": ["PPTX Fixture"],
                    "source_anchor": {
                        "slide_no": 1,
                        "shape_index": 2,
                        "text_shape_index": 2,
                        "slide_text_shape_index": 2,
                    },
                    "slide_no": 1,
                    "metadata": {"source": "pptx", "is_title": False},
                },
                {
                    "block_seq": 2,
                    "block_type": "table",
                    "parent_block_seq": 0,
                    "content_text": "Metric\tValue\nQuality\tBaseline",
                    "content_markdown": (
                        "| Metric | Value |\n" "| --- | --- |\n" "| Quality | Baseline |"
                    ),
                    "heading_path": ["PPTX Fixture"],
                    "source_anchor": {
                        "slide_no": 1,
                        "shape_index": 3,
                        "table_index": 1,
                        "slide_table_index": 1,
                    },
                    "slide_no": 1,
                    "metadata": {
                        "source": "pptx",
                        "row_count": 2,
                        "column_count": 2,
                    },
                },
            ],
        },
    ),
    ExtractionFixtureCase(
        case_id="xlsx",
        file_name="fixture.xlsx",
        profile_name=LOCAL_XLSX_PROFILE_NAME,
        detected_file_type="xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        write_source=write_xlsx_fixture,
        expected_snapshot={
            "status": "succeeded",
            "runtime_metadata": {
                "profile_name": LOCAL_XLSX_PROFILE_NAME,
                "extractor_name": "local_xlsx",
                "extractor_version": "0.1.0",
                "library": "openpyxl",
                "sheet_count": 1,
                "extracted_sheet_count": 1,
                "table_count": 1,
                "cell_count": 6,
                "block_count": 2,
                "preserve_sheet_boundaries": True,
                "emit_markdown_tables": True,
            },
            "artifact": {
                "artifact_type": "normalized_markdown",
                "content_text": (
                    "<!-- sheet: Fixture Sheet -->\n\n"
                    "# Fixture Sheet\n\n"
                    "| Metric | Value |\n"
                    "| --- | --- |\n"
                    "| Accuracy | High |\n"
                    "| Latency | Low |"
                ),
                "metadata": {
                    "parser_name": "local_xlsx",
                    "parser_version": "0.1.0",
                    "library": "openpyxl",
                    "sheet_count": 1,
                    "extracted_sheet_count": 1,
                    "table_count": 1,
                    "cell_count": 6,
                    "block_count": 2,
                    "preserve_sheet_boundaries": True,
                    "emit_markdown_tables": True,
                    "formulas_resolved_from_cached_values": True,
                },
            },
            "blocks": [
                {
                    "block_seq": 0,
                    "block_type": "heading",
                    "content_text": "Fixture Sheet",
                    "content_markdown": "# Fixture Sheet",
                    "heading_path": ["Fixture Sheet"],
                    "source_anchor": {
                        "sheet_index": 1,
                        "sheet_name": "Fixture Sheet",
                    },
                    "sheet_name": "Fixture Sheet",
                    "metadata": {"source": "xlsx", "sheet_index": 1},
                },
                {
                    "block_seq": 1,
                    "block_type": "table",
                    "parent_block_seq": 0,
                    "content_text": ("Metric\tValue\n" "Accuracy\tHigh\n" "Latency\tLow"),
                    "content_markdown": (
                        "| Metric | Value |\n"
                        "| --- | --- |\n"
                        "| Accuracy | High |\n"
                        "| Latency | Low |"
                    ),
                    "heading_path": ["Fixture Sheet"],
                    "source_anchor": {
                        "sheet_index": 1,
                        "sheet_name": "Fixture Sheet",
                        "cell_range": "A1:B3",
                        "table_index": 1,
                    },
                    "sheet_name": "Fixture Sheet",
                    "cell_range": "A1:B3",
                    "metadata": {
                        "source": "xlsx",
                        "row_count": 3,
                        "column_count": 2,
                        "cell_count": 6,
                    },
                },
            ],
        },
    ),
    ExtractionFixtureCase(
        case_id="hwpx",
        file_name="fixture.hwpx",
        profile_name=LOCAL_HWPX_PROFILE_NAME,
        detected_file_type="hwpx",
        mime_type="application/vnd.hancom.hwpx",
        write_source=write_hwpx_fixture,
        expected_snapshot={
            "status": "succeeded",
            "runtime_metadata": {
                "profile_name": LOCAL_HWPX_PROFILE_NAME,
                "extractor_name": "local_hwpx",
                "extractor_version": "0.1.0",
                "container_format": "hwpx_zip_xml",
                "section_count": 1,
                "extracted_section_count": 1,
                "paragraph_count": 1,
                "table_count": 1,
                "block_count": 2,
                "preserve_sections": True,
                "preserve_tables": True,
            },
            "artifact": {
                "artifact_type": "normalized_markdown",
                "content_text": (
                    "<!-- section: 1 -->\n\n"
                    "HWPX body paragraph.\n\n"
                    "| Metric | Value |\n"
                    "| --- | --- |\n"
                    "| Quality | Baseline |"
                ),
                "metadata": {
                    "parser_name": "local_hwpx",
                    "parser_version": "0.1.0",
                    "container_format": "hwpx_zip_xml",
                    "section_count": 1,
                    "extracted_section_count": 1,
                    "paragraph_count": 1,
                    "table_count": 1,
                    "block_count": 2,
                    "preserve_sections": True,
                    "preserve_tables": True,
                },
            },
            "blocks": [
                {
                    "block_seq": 0,
                    "block_type": "paragraph",
                    "content_text": "HWPX body paragraph.",
                    "content_markdown": "HWPX body paragraph.",
                    "heading_path": ["Section 1"],
                    "source_anchor": {
                        "section_index": 1,
                        "section_name": "Contents/section0.xml",
                        "paragraph_index": 1,
                        "section_paragraph_index": 1,
                    },
                    "metadata": {
                        "source": "hwpx",
                        "section_name": "Contents/section0.xml",
                    },
                },
                {
                    "block_seq": 1,
                    "block_type": "table",
                    "content_text": "Metric\tValue\nQuality\tBaseline",
                    "content_markdown": (
                        "| Metric | Value |\n" "| --- | --- |\n" "| Quality | Baseline |"
                    ),
                    "heading_path": ["Section 1"],
                    "source_anchor": {
                        "section_index": 1,
                        "section_name": "Contents/section0.xml",
                        "table_index": 1,
                        "section_table_index": 1,
                    },
                    "metadata": {
                        "source": "hwpx",
                        "section_name": "Contents/section0.xml",
                        "row_count": 2,
                        "column_count": 2,
                    },
                },
            ],
        },
    ),
)
