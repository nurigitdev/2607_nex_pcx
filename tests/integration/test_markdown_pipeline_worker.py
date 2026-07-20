from collections import Counter
from hashlib import sha256
from io import BytesIO
from pathlib import Path
from uuid import uuid4
from zipfile import ZipFile

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.core.bm25_keyword_index import (
    list_chunk_keyword_statistics,
    list_chunk_keyword_terms,
    refresh_chunk_policy_keyword_index,
)
from app.core.chunks import DEFAULT_CHUNK_POLICY_NAME, list_document_chunks
from app.core.database import connect, fetch_one
from app.core.embedding_jobs import list_embedding_jobs
from app.core.file_uploads import store_upload
from app.core.ingestion_artifacts import (
    list_document_blocks,
    list_document_extraction_artifacts,
    list_document_extraction_runs,
)
from app.core.pipeline_jobs import list_pipeline_job_events
from app.core.pipeline_worker import (
    DEFAULT_PIPELINE_CHUNK_POLICY_NAMES,
    process_next_markdown_pipeline_job,
)

pytestmark = pytest.mark.integration


def cleanup_checksum(database_url: str, checksum: str) -> None:
    with connect(database_url) as connection:
        with connection.cursor() as cursor:
            cursor.execute("DELETE FROM files WHERE sha256_checksum = %s", (checksum,))
    for policy_name in DEFAULT_PIPELINE_CHUNK_POLICY_NAMES:
        refresh_chunk_policy_keyword_index(database_url, chunk_policy_name=policy_name)


def prioritize_job(database_url: str, job_id: int) -> None:
    with connect(database_url) as connection:
        with connection.cursor() as cursor:
            cursor.execute(
                "UPDATE pipeline_jobs SET priority = 0 WHERE job_id = %s",
                (job_id,),
            )


def assert_policy_chunk_counts(chunks, expected_chunks_per_policy: int) -> None:
    assert Counter(chunk.chunk_policy_name for chunk in chunks) == {
        policy_name: expected_chunks_per_policy
        for policy_name in DEFAULT_PIPELINE_CHUNK_POLICY_NAMES
    }


def make_minimal_pdf(lines: list[str]) -> bytes:
    text_operations: list[str] = []
    y_position = 760
    for line in lines:
        escaped = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        text_operations.append(f"BT /F1 12 Tf 72 {y_position} Td ({escaped}) Tj ET")
        y_position -= 18

    stream = "\n".join(text_operations).encode("ascii")
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        (
            b"3 0 obj\n"
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>\n"
            b"endobj\n"
        ),
        b"4 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
        (
            b"5 0 obj\n<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"\nendstream\nendobj\n"
        ),
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for pdf_object in objects:
        offsets.append(len(output))
        output.extend(pdf_object)

    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        (
            f"trailer\n<< /Root 1 0 R /Size {len(objects) + 1} >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(output)


def make_sample_docx_bytes(unique_text: str) -> bytes:
    buffer = BytesIO()
    document = Document()
    document.add_heading("Slice 223 DOCX", level=1)
    document.add_paragraph(f"DOCX pipeline integration {unique_text}.")
    document.add_heading("Measurements", level=2)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Quality"
    table.cell(1, 1).text = "Baseline"
    document.save(buffer)
    return buffer.getvalue()


def make_sample_pptx_bytes(unique_text: str) -> bytes:
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slice 224 PPTX"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    textbox.text_frame.text = f"PPTX pipeline integration {unique_text}."

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(6), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Quality"
    table.cell(1, 1).text = "Baseline"
    presentation.save(buffer)
    return buffer.getvalue()


def make_sample_xlsx_bytes(unique_text: str) -> bytes:
    buffer = BytesIO()
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Slice 225 Metrics"
    worksheet.append(["Metric", "Value"])
    worksheet.append(["Quality", "Baseline"])
    worksheet.append(["Trace", unique_text])
    workbook.save(buffer)
    return buffer.getvalue()


def make_sample_hwpx_bytes(unique_text: str) -> bytes:
    buffer = BytesIO()
    section_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p>
    <hp:run><hp:t>HWPX pipeline integration {unique_text}.</hp:t></hp:run>
  </hp:p>
  <hp:tbl>
    <hp:tr>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Metric</hp:t></hp:run></hp:p></hp:subList></hp:tc>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Value</hp:t></hp:run></hp:p></hp:subList></hp:tc>
    </hp:tr>
    <hp:tr>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Quality</hp:t></hp:run></hp:p></hp:subList></hp:tc>
      <hp:tc><hp:subList><hp:p><hp:run><hp:t>Baseline</hp:t></hp:run></hp:p></hp:subList></hp:tc>
    </hp:tr>
  </hp:tbl>
</hp:sec>
"""
    with ZipFile(buffer, "w") as archive:
        archive.writestr("mimetype", "application/hwp+zip")
        archive.writestr("Contents/section0.xml", section_xml)
    return buffer.getvalue()


def test_markdown_pipeline_worker_parses_and_stores_chunks(
    migrated_database_url: str,
    tmp_path: Path,
) -> None:
    unique_text = uuid4().hex
    content = f"""# Slice 017

Markdown pipeline worker integration test {unique_text}.

## Details

This section should become another heading-aware chunk.
""".encode()
    checksum = sha256(content).hexdigest()

    try:
        upload_result = store_upload(
            database_url=migrated_database_url,
            upload_stream=BytesIO(content),
            original_file_name="slice-017-worker.md",
            storage_dir=tmp_path,
            mime_type="text/markdown",
            document_group="slice-017",
            security_level="internal",
            uploaded_by="integration-test",
        )
        assert upload_result.pipeline_job is not None
        prioritize_job(migrated_database_url, upload_result.pipeline_job.job_id)

        result = process_next_markdown_pipeline_job(
            migrated_database_url,
            worker_name="slice-017-test-worker",
        )

        assert result.processed is True
        assert result.job is not None
        assert result.job.job_id == upload_result.pipeline_job.job_id
        assert result.job.status == "succeeded"
        assert result.job.stage == "completed"
        assert result.job.progress_percent == 100
        assert result.chunk_count == 6
        assert result.embedding_job_count == 24
        assert result.bm25_term_row_count > 0
        assert result.bm25_statistics_row_count > 0
        assert {
            policy_result.chunk_policy_name: policy_result.chunk_count
            for policy_result in result.policy_results
        } == {policy_name: 2 for policy_name in DEFAULT_PIPELINE_CHUNK_POLICY_NAMES}
        assert all(policy_result.bm25_term_row_count > 0 for policy_result in result.policy_results)
        assert all(
            policy_result.bm25_statistics_row_count > 0 for policy_result in result.policy_results
        )

        chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
        )
        default_chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )
        extraction_runs = list_document_extraction_runs(
            migrated_database_url,
            upload_result.file.document_id,
        )
        extraction_artifacts = list_document_extraction_artifacts(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_type="normalized_markdown",
        )
        document_blocks = list_document_blocks(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_id=extraction_artifacts[0].artifact_id,
        )
        embedding_jobs = [
            job
            for chunk in chunks
            for job in list_embedding_jobs(migrated_database_url, chunk_id=chunk.chunk_id)
        ]
        default_chunk_terms = list_chunk_keyword_terms(
            migrated_database_url,
            chunk_id=default_chunks[0].chunk_id,
        )
        default_policy_statistics = list_chunk_keyword_statistics(
            migrated_database_url,
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )
        file_row = fetch_one(
            migrated_database_url,
            """
            SELECT
                parser_name,
                parser_version,
                parse_status,
                parse_error_message,
                extracted_text_size
            FROM files
            WHERE file_id = %s
            """,
            (upload_result.file.file_id,),
        )
        events = list_pipeline_job_events(
            migrated_database_url,
            upload_result.pipeline_job.job_id,
        )

        assert_policy_chunk_counts(chunks, 2)
        assert [chunk.chunk_seq for chunk in default_chunks] == [0, 1]
        assert default_chunks[0].chunk_text.startswith("# Slice 017")
        assert default_chunks[1].heading_path == ("Slice 017", "Details")
        assert extraction_runs[0].status == "succeeded"
        assert extraction_runs[0].extraction_profile_name == "local_markdown_default"
        assert extraction_artifacts[0].artifact_type == "normalized_markdown"
        assert [block.block_type for block in document_blocks] == [
            "heading",
            "paragraph",
            "heading",
            "paragraph",
        ]
        assert default_chunks[0].artifact_id == extraction_artifacts[0].artifact_id
        assert default_chunks[0].block_id == document_blocks[0].block_id
        assert default_chunks[0].metadata["block_ids"] == [
            document_blocks[0].block_id,
            document_blocks[1].block_id,
        ]
        assert len(embedding_jobs) == 24
        assert {job.profile_name for job in embedding_jobs} == {
            "kure_v1_1024",
            "bge_m3_1024",
            "qwen3_4b_1000",
            "qwen3_4b_2560",
        }
        assert {job.status for job in embedding_jobs} == {"pending"}
        assert {term.term for term in default_chunk_terms} >= {"slice", "017"}
        assert default_policy_statistics
        assert all(
            statistic.corpus_chunk_count >= len(default_chunks)
            for statistic in default_policy_statistics
        )
        assert file_row["parser_name"] == "markdown"
        assert file_row["parser_version"] == "0.1.0"
        assert file_row["parse_status"] == "succeeded"
        assert file_row["parse_error_message"] is None
        assert file_row["extracted_text_size"] > 0
        assert {event.event_type for event in events} >= {
            "created",
            "claimed",
            "progress",
            "stage_succeeded",
        }
    finally:
        cleanup_checksum(migrated_database_url, checksum)


def test_markdown_pipeline_worker_can_limit_processing_to_one_chunk_policy(
    migrated_database_url: str,
    tmp_path: Path,
) -> None:
    unique_text = uuid4().hex
    content = f"""# Single Policy

Markdown pipeline worker single policy integration test {unique_text}.
""".encode()
    checksum = sha256(content).hexdigest()

    try:
        upload_result = store_upload(
            database_url=migrated_database_url,
            upload_stream=BytesIO(content),
            original_file_name="slice-259-single-policy-worker.md",
            storage_dir=tmp_path,
            mime_type="text/markdown",
            document_group="slice-259",
            security_level="internal",
            uploaded_by="integration-test",
        )
        assert upload_result.pipeline_job is not None
        prioritize_job(migrated_database_url, upload_result.pipeline_job.job_id)

        result = process_next_markdown_pipeline_job(
            migrated_database_url,
            worker_name="slice-259-single-policy-worker",
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )

        chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
        )
        embedding_jobs = [
            job
            for chunk in chunks
            for job in list_embedding_jobs(migrated_database_url, chunk_id=chunk.chunk_id)
        ]

        assert result.processed is True
        assert result.job is not None
        assert result.job.status == "succeeded"
        assert result.chunk_count == 1
        assert result.embedding_job_count == 4
        assert result.bm25_term_row_count > 0
        assert result.bm25_statistics_row_count > 0
        assert result.policy_results[0].chunk_policy_name == DEFAULT_CHUNK_POLICY_NAME
        assert result.policy_results[0].bm25_term_row_count > 0
        assert result.policy_results[0].bm25_statistics_row_count > 0
        assert {chunk.chunk_policy_name for chunk in chunks} == {DEFAULT_CHUNK_POLICY_NAME}
        assert len(embedding_jobs) == 4
    finally:
        cleanup_checksum(migrated_database_url, checksum)


def test_markdown_pipeline_worker_extracts_docx_paragraphs_and_tables(
    migrated_database_url: str,
    tmp_path: Path,
) -> None:
    unique_text = uuid4().hex
    content = make_sample_docx_bytes(unique_text)
    checksum = sha256(content).hexdigest()

    try:
        upload_result = store_upload(
            database_url=migrated_database_url,
            upload_stream=BytesIO(content),
            original_file_name="slice-223-worker.docx",
            storage_dir=tmp_path,
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            document_group="slice-223",
            security_level="internal",
            uploaded_by="integration-test",
        )
        assert upload_result.pipeline_job is not None
        prioritize_job(migrated_database_url, upload_result.pipeline_job.job_id)

        result = process_next_markdown_pipeline_job(
            migrated_database_url,
            worker_name="slice-223-test-worker",
        )

        assert result.processed is True
        assert result.job is not None
        assert result.job.job_id == upload_result.pipeline_job.job_id
        assert result.job.status == "succeeded"
        assert result.chunk_count == 6
        assert result.embedding_job_count == 24

        chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
        )
        default_chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )
        extraction_runs = list_document_extraction_runs(
            migrated_database_url,
            upload_result.file.document_id,
        )
        extraction_artifacts = list_document_extraction_artifacts(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_type="normalized_markdown",
        )
        document_blocks = list_document_blocks(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_id=extraction_artifacts[0].artifact_id,
        )
        embedding_jobs = [
            job
            for chunk in chunks
            for job in list_embedding_jobs(migrated_database_url, chunk_id=chunk.chunk_id)
        ]
        file_row = fetch_one(
            migrated_database_url,
            """
            SELECT
                parser_name,
                parser_version,
                parse_status,
                parse_error_message,
                extracted_text_size
            FROM files
            WHERE file_id = %s
            """,
            (upload_result.file.file_id,),
        )

        assert extraction_runs[0].status == "succeeded"
        assert extraction_runs[0].extraction_profile_name == "local_docx_default"
        assert extraction_runs[0].runtime_metadata["library"] == "python-docx"
        assert extraction_artifacts[0].metadata["preserve_tables"] is True
        assert extraction_artifacts[0].metadata["heading_count"] == 2
        assert [block.block_type for block in document_blocks] == [
            "heading",
            "paragraph",
            "heading",
            "table",
        ]
        assert document_blocks[0].content_markdown == "# Slice 223 DOCX"
        assert document_blocks[3].metadata["source"] == "docx"
        assert_policy_chunk_counts(chunks, 2)
        assert default_chunks[0].chunk_text.startswith("# Slice 223 DOCX")
        assert default_chunks[1].heading_path == ("Slice 223 DOCX", "Measurements")
        assert default_chunks[0].artifact_id == extraction_artifacts[0].artifact_id
        assert default_chunks[0].block_id == document_blocks[0].block_id
        assert len(embedding_jobs) == 24
        assert {job.status for job in embedding_jobs} == {"pending"}
        assert file_row["parser_name"] == "local_docx"
        assert file_row["parser_version"] == "0.1.0"
        assert file_row["parse_status"] == "succeeded"
        assert file_row["parse_error_message"] is None
        assert file_row["extracted_text_size"] > 0
    finally:
        cleanup_checksum(migrated_database_url, checksum)


def test_markdown_pipeline_worker_extracts_pptx_slides_and_tables(
    migrated_database_url: str,
    tmp_path: Path,
) -> None:
    unique_text = uuid4().hex
    content = make_sample_pptx_bytes(unique_text)
    checksum = sha256(content).hexdigest()

    try:
        upload_result = store_upload(
            database_url=migrated_database_url,
            upload_stream=BytesIO(content),
            original_file_name="slice-224-worker.pptx",
            storage_dir=tmp_path,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            document_group="slice-224",
            security_level="internal",
            uploaded_by="integration-test",
        )
        assert upload_result.pipeline_job is not None
        prioritize_job(migrated_database_url, upload_result.pipeline_job.job_id)

        result = process_next_markdown_pipeline_job(
            migrated_database_url,
            worker_name="slice-224-test-worker",
        )

        assert result.processed is True
        assert result.job is not None
        assert result.job.job_id == upload_result.pipeline_job.job_id
        assert result.job.status == "succeeded"
        assert result.chunk_count == 3
        assert result.embedding_job_count == 12

        chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
        )
        default_chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )
        extraction_runs = list_document_extraction_runs(
            migrated_database_url,
            upload_result.file.document_id,
        )
        extraction_artifacts = list_document_extraction_artifacts(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_type="normalized_markdown",
        )
        document_blocks = list_document_blocks(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_id=extraction_artifacts[0].artifact_id,
        )
        embedding_jobs = [
            job
            for chunk in chunks
            for job in list_embedding_jobs(migrated_database_url, chunk_id=chunk.chunk_id)
        ]
        file_row = fetch_one(
            migrated_database_url,
            """
            SELECT
                parser_name,
                parser_version,
                parse_status,
                parse_error_message,
                extracted_text_size
            FROM files
            WHERE file_id = %s
            """,
            (upload_result.file.file_id,),
        )

        assert extraction_runs[0].status == "succeeded"
        assert extraction_runs[0].extraction_profile_name == "local_pptx_default"
        assert extraction_runs[0].runtime_metadata["library"] == "python-pptx"
        assert extraction_artifacts[0].metadata["preserve_slide_boundaries"] is True
        assert extraction_artifacts[0].metadata["preserve_tables"] is True
        assert extraction_artifacts[0].metadata["slide_count"] == 1
        assert extraction_artifacts[0].metadata["table_count"] == 1
        assert [block.block_type for block in document_blocks] == [
            "heading",
            "paragraph",
            "table",
        ]
        assert [block.slide_no for block in document_blocks] == [1, 1, 1]
        assert document_blocks[0].content_markdown == "# Slice 224 PPTX"
        assert document_blocks[2].metadata["source"] == "pptx"
        assert_policy_chunk_counts(chunks, 1)
        assert default_chunks[0].chunk_text.startswith("# Slice 224 PPTX")
        assert default_chunks[0].artifact_id == extraction_artifacts[0].artifact_id
        assert default_chunks[0].block_id == document_blocks[0].block_id
        assert len(embedding_jobs) == 12
        assert {job.status for job in embedding_jobs} == {"pending"}
        assert file_row["parser_name"] == "local_pptx"
        assert file_row["parser_version"] == "0.1.0"
        assert file_row["parse_status"] == "succeeded"
        assert file_row["parse_error_message"] is None
        assert file_row["extracted_text_size"] > 0
    finally:
        cleanup_checksum(migrated_database_url, checksum)


def test_markdown_pipeline_worker_extracts_xlsx_sheets_and_tables(
    migrated_database_url: str,
    tmp_path: Path,
) -> None:
    unique_text = uuid4().hex
    content = make_sample_xlsx_bytes(unique_text)
    checksum = sha256(content).hexdigest()

    try:
        upload_result = store_upload(
            database_url=migrated_database_url,
            upload_stream=BytesIO(content),
            original_file_name="slice-225-worker.xlsx",
            storage_dir=tmp_path,
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            document_group="slice-225",
            security_level="internal",
            uploaded_by="integration-test",
        )
        assert upload_result.pipeline_job is not None
        prioritize_job(migrated_database_url, upload_result.pipeline_job.job_id)

        result = process_next_markdown_pipeline_job(
            migrated_database_url,
            worker_name="slice-225-test-worker",
        )

        assert result.processed is True
        assert result.job is not None
        assert result.job.job_id == upload_result.pipeline_job.job_id
        assert result.job.status == "succeeded"
        assert result.chunk_count == 3
        assert result.embedding_job_count == 12

        chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
        )
        default_chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )
        extraction_runs = list_document_extraction_runs(
            migrated_database_url,
            upload_result.file.document_id,
        )
        extraction_artifacts = list_document_extraction_artifacts(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_type="normalized_markdown",
        )
        document_blocks = list_document_blocks(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_id=extraction_artifacts[0].artifact_id,
        )
        embedding_jobs = [
            job
            for chunk in chunks
            for job in list_embedding_jobs(migrated_database_url, chunk_id=chunk.chunk_id)
        ]
        file_row = fetch_one(
            migrated_database_url,
            """
            SELECT
                parser_name,
                parser_version,
                parse_status,
                parse_error_message,
                extracted_text_size
            FROM files
            WHERE file_id = %s
            """,
            (upload_result.file.file_id,),
        )

        assert extraction_runs[0].status == "succeeded"
        assert extraction_runs[0].extraction_profile_name == "local_xlsx_default"
        assert extraction_runs[0].runtime_metadata["library"] == "openpyxl"
        assert extraction_artifacts[0].metadata["preserve_sheet_boundaries"] is True
        assert extraction_artifacts[0].metadata["emit_markdown_tables"] is True
        assert extraction_artifacts[0].metadata["sheet_count"] == 1
        assert extraction_artifacts[0].metadata["table_count"] == 1
        assert [block.block_type for block in document_blocks] == ["heading", "table"]
        assert [block.sheet_name for block in document_blocks] == [
            "Slice 225 Metrics",
            "Slice 225 Metrics",
        ]
        assert document_blocks[0].content_markdown == "# Slice 225 Metrics"
        assert document_blocks[1].metadata["source"] == "xlsx"
        assert document_blocks[1].cell_range == "A1:B3"
        assert_policy_chunk_counts(chunks, 1)
        assert default_chunks[0].chunk_text.startswith("# Slice 225 Metrics")
        assert default_chunks[0].artifact_id == extraction_artifacts[0].artifact_id
        assert default_chunks[0].block_id == document_blocks[0].block_id
        assert len(embedding_jobs) == 12
        assert {job.status for job in embedding_jobs} == {"pending"}
        assert file_row["parser_name"] == "local_xlsx"
        assert file_row["parser_version"] == "0.1.0"
        assert file_row["parse_status"] == "succeeded"
        assert file_row["parse_error_message"] is None
        assert file_row["extracted_text_size"] > 0
    finally:
        cleanup_checksum(migrated_database_url, checksum)


def test_markdown_pipeline_worker_extracts_hwpx_paragraphs_and_tables(
    migrated_database_url: str,
    tmp_path: Path,
) -> None:
    unique_text = uuid4().hex
    content = make_sample_hwpx_bytes(unique_text)
    checksum = sha256(content).hexdigest()

    try:
        upload_result = store_upload(
            database_url=migrated_database_url,
            upload_stream=BytesIO(content),
            original_file_name="slice-226-worker.hwpx",
            storage_dir=tmp_path,
            mime_type="application/vnd.hancom.hwpx",
            document_group="slice-226",
            security_level="internal",
            uploaded_by="integration-test",
        )
        assert upload_result.pipeline_job is not None
        prioritize_job(migrated_database_url, upload_result.pipeline_job.job_id)

        result = process_next_markdown_pipeline_job(
            migrated_database_url,
            worker_name="slice-226-test-worker",
        )

        assert result.processed is True
        assert result.job is not None
        assert result.job.job_id == upload_result.pipeline_job.job_id
        assert result.job.status == "succeeded"
        assert result.chunk_count == 3
        assert result.embedding_job_count == 12

        chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
        )
        default_chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )
        extraction_runs = list_document_extraction_runs(
            migrated_database_url,
            upload_result.file.document_id,
        )
        extraction_artifacts = list_document_extraction_artifacts(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_type="normalized_markdown",
        )
        document_blocks = list_document_blocks(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_id=extraction_artifacts[0].artifact_id,
        )
        embedding_jobs = [
            job
            for chunk in chunks
            for job in list_embedding_jobs(migrated_database_url, chunk_id=chunk.chunk_id)
        ]
        file_row = fetch_one(
            migrated_database_url,
            """
            SELECT
                parser_name,
                parser_version,
                parse_status,
                parse_error_message,
                extracted_text_size
            FROM files
            WHERE file_id = %s
            """,
            (upload_result.file.file_id,),
        )

        assert extraction_runs[0].status == "succeeded"
        assert extraction_runs[0].extraction_profile_name == "local_hwpx_default"
        assert extraction_runs[0].runtime_metadata["container_format"] == "hwpx_zip_xml"
        assert extraction_artifacts[0].metadata["preserve_sections"] is True
        assert extraction_artifacts[0].metadata["preserve_tables"] is True
        assert extraction_artifacts[0].metadata["section_count"] == 1
        assert extraction_artifacts[0].metadata["table_count"] == 1
        assert [block.block_type for block in document_blocks] == ["paragraph", "table"]
        assert document_blocks[0].content_markdown.startswith("HWPX pipeline integration")
        assert document_blocks[1].metadata["source"] == "hwpx"
        assert_policy_chunk_counts(chunks, 1)
        assert default_chunks[0].chunk_text.startswith("HWPX pipeline integration")
        assert default_chunks[0].artifact_id == extraction_artifacts[0].artifact_id
        assert default_chunks[0].block_id == document_blocks[0].block_id
        assert len(embedding_jobs) == 12
        assert {job.status for job in embedding_jobs} == {"pending"}
        assert file_row["parser_name"] == "local_hwpx"
        assert file_row["parser_version"] == "0.1.0"
        assert file_row["parse_status"] == "succeeded"
        assert file_row["parse_error_message"] is None
        assert file_row["extracted_text_size"] > 0
    finally:
        cleanup_checksum(migrated_database_url, checksum)


def test_markdown_pipeline_worker_extracts_pdf_text_layer(
    migrated_database_url: str,
    tmp_path: Path,
) -> None:
    content = make_minimal_pdf(
        [
            "Slice 222 PDF",
            "",
            f"PDF text layer integration {uuid4().hex}.",
            "This content should become a PDF document block.",
        ]
    )
    checksum = sha256(content).hexdigest()

    try:
        upload_result = store_upload(
            database_url=migrated_database_url,
            upload_stream=BytesIO(content),
            original_file_name="slice-222-worker.pdf",
            storage_dir=tmp_path,
            mime_type="application/pdf",
            document_group="slice-222",
            security_level="internal",
            uploaded_by="integration-test",
        )
        assert upload_result.pipeline_job is not None
        prioritize_job(migrated_database_url, upload_result.pipeline_job.job_id)

        result = process_next_markdown_pipeline_job(
            migrated_database_url,
            worker_name="slice-222-test-worker",
        )

        assert result.processed is True
        assert result.job is not None
        assert result.job.job_id == upload_result.pipeline_job.job_id
        assert result.job.status == "succeeded"
        assert result.chunk_count == 3
        assert result.embedding_job_count == 12

        chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
        )
        default_chunks = list_document_chunks(
            migrated_database_url,
            upload_result.file.document_id,
            chunk_policy_name=DEFAULT_CHUNK_POLICY_NAME,
        )
        extraction_runs = list_document_extraction_runs(
            migrated_database_url,
            upload_result.file.document_id,
        )
        extraction_artifacts = list_document_extraction_artifacts(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_type="normalized_markdown",
        )
        document_blocks = list_document_blocks(
            migrated_database_url,
            upload_result.file.document_id,
            artifact_id=extraction_artifacts[0].artifact_id,
        )
        embedding_jobs = [
            job
            for chunk in chunks
            for job in list_embedding_jobs(migrated_database_url, chunk_id=chunk.chunk_id)
        ]
        file_row = fetch_one(
            migrated_database_url,
            """
            SELECT
                parser_name,
                parser_version,
                parse_status,
                parse_error_message,
                extracted_text_size
            FROM files
            WHERE file_id = %s
            """,
            (upload_result.file.file_id,),
        )

        assert extraction_runs[0].status == "succeeded"
        assert extraction_runs[0].extraction_profile_name == "local_pdf_text_default"
        assert extraction_runs[0].runtime_metadata["library"] == "pypdf"
        assert extraction_artifacts[0].metadata["text_layer_only"] is True
        assert extraction_artifacts[0].metadata["ocr_enabled"] is False
        assert [block.page_no for block in document_blocks] == [1]
        assert document_blocks[0].metadata["source"] == "pdf_text_layer"
        assert_policy_chunk_counts(chunks, 1)
        assert default_chunks[0].chunk_text.startswith("Slice 222 PDF")
        assert default_chunks[0].artifact_id == extraction_artifacts[0].artifact_id
        assert default_chunks[0].block_id == document_blocks[0].block_id
        assert len(embedding_jobs) == 12
        assert {job.status for job in embedding_jobs} == {"pending"}
        assert file_row["parser_name"] == "local_pdf_text"
        assert file_row["parser_version"] == "0.1.0"
        assert file_row["parse_status"] == "succeeded"
        assert file_row["parse_error_message"] is None
        assert file_row["extracted_text_size"] > 0
    finally:
        cleanup_checksum(migrated_database_url, checksum)
