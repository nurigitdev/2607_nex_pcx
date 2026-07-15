"""Local extraction runtime registry for local document sources."""

from collections.abc import Callable
from dataclasses import dataclass
from hashlib import sha256
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path
from time import perf_counter
from typing import Any

from psycopg import Connection

from app.core.chunking import count_chunk_tokens
from app.core.database import connect
from app.core.document_parsers import (
    PARSER_NAME_MARKDOWN,
    PARSER_VERSION_MARKDOWN,
    MarkdownParser,
    ParsedBlock,
)
from app.core.extraction_runtime import (
    ExtractionRuntimeArtifact,
    ExtractionRuntimeBlock,
    ExtractionRuntimeRequest,
    ExtractionRuntimeResult,
    validate_runtime_request,
    validate_runtime_result,
)
from app.core.ingestion_artifacts import (
    DocumentBlockInput,
    DocumentBlockRecord,
    ExtractionArtifactInput,
    ExtractionArtifactRecord,
    ExtractionRunInput,
    ExtractionRunRecord,
    InvalidIngestionArtifactError,
    create_document_block_in_connection,
    create_extraction_artifact_in_connection,
    create_extraction_run_in_connection,
    get_extraction_profile_in_connection,
)

LOCAL_MARKDOWN_PROFILE_NAME = "local_markdown_default"
LOCAL_PLAIN_TEXT_PROFILE_NAME = "local_plain_text_default"
LOCAL_PDF_TEXT_PROFILE_NAME = "local_pdf_text_default"
LOCAL_DOCX_PROFILE_NAME = "local_docx_default"
LOCAL_PPTX_PROFILE_NAME = "local_pptx_default"
LOCAL_XLSX_PROFILE_NAME = "local_xlsx_default"
LOCAL_PLAIN_TEXT_EXTRACTOR_NAME = "local_plain_text"
LOCAL_PLAIN_TEXT_EXTRACTOR_VERSION = "0.1.0"
LOCAL_PDF_TEXT_EXTRACTOR_NAME = "local_pdf_text"
LOCAL_PDF_TEXT_EXTRACTOR_VERSION = "0.1.0"
LOCAL_DOCX_EXTRACTOR_NAME = "local_docx"
LOCAL_DOCX_EXTRACTOR_VERSION = "0.1.0"
LOCAL_PPTX_EXTRACTOR_NAME = "local_pptx"
LOCAL_PPTX_EXTRACTOR_VERSION = "0.1.0"
LOCAL_XLSX_EXTRACTOR_NAME = "local_xlsx"
LOCAL_XLSX_EXTRACTOR_VERSION = "0.1.0"
PDF_TEXT_LIBRARY_NAME = "pypdf"
DOCX_LIBRARY_NAME = "python-docx"
PPTX_LIBRARY_NAME = "python-pptx"
XLSX_LIBRARY_NAME = "openpyxl"

ERROR_CODE_LOCAL_SOURCE_NOT_FOUND = "LOCAL_SOURCE_NOT_FOUND"
ERROR_CODE_LOCAL_SOURCE_EMPTY = "LOCAL_SOURCE_EMPTY"
ERROR_CODE_LOCAL_UNSUPPORTED_PROFILE = "LOCAL_UNSUPPORTED_PROFILE"
ERROR_CODE_LOCAL_UNSUPPORTED_FILE_TYPE = "LOCAL_UNSUPPORTED_FILE_TYPE"
ERROR_CODE_LOCAL_TEXT_DECODE_FAILED = "LOCAL_TEXT_DECODE_FAILED"
ERROR_CODE_LOCAL_PDF_READ_FAILED = "LOCAL_PDF_READ_FAILED"
ERROR_CODE_LOCAL_PDF_TEXT_LAYER_EMPTY = "LOCAL_PDF_TEXT_LAYER_EMPTY"
ERROR_CODE_LOCAL_DOCX_READ_FAILED = "LOCAL_DOCX_READ_FAILED"
ERROR_CODE_LOCAL_DOCX_EMPTY = "LOCAL_DOCX_EMPTY"
ERROR_CODE_LOCAL_PPTX_READ_FAILED = "LOCAL_PPTX_READ_FAILED"
ERROR_CODE_LOCAL_PPTX_EMPTY = "LOCAL_PPTX_EMPTY"
ERROR_CODE_LOCAL_XLSX_READ_FAILED = "LOCAL_XLSX_READ_FAILED"
ERROR_CODE_LOCAL_XLSX_EMPTY = "LOCAL_XLSX_EMPTY"


LocalExtractionCallable = Callable[
    [ExtractionRuntimeRequest, Path, float],
    ExtractionRuntimeResult,
]


@dataclass(frozen=True)
class LocalExtractionHandler:
    profile_name: str
    parser_name: str
    parser_version: str
    extractor_name: str
    extractor_version: str
    supported_file_types: tuple[str, ...]
    extract: LocalExtractionCallable

    def supports_file_type(self, file_type: str | None) -> bool:
        normalized = normalize_file_type(file_type)
        return bool(normalized) and normalized in self.supported_file_types


@dataclass(frozen=True)
class PersistedExtractionRuntimeResult:
    run: ExtractionRunRecord
    artifacts: tuple[ExtractionArtifactRecord, ...] = ()
    blocks: tuple[DocumentBlockRecord, ...] = ()


def run_local_extraction(request: ExtractionRuntimeRequest) -> ExtractionRuntimeResult:
    """Run the registered local extraction profile and return a terminal runtime result."""

    validate_runtime_request(request)
    started = perf_counter()
    source_path = Path(request.storage_path)
    profile_name = request.extraction_profile_name.strip()
    file_type = normalize_file_type(request.detected_file_type) or normalize_file_type(
        source_path.suffix
    )

    if not source_path.is_file():
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_SOURCE_NOT_FOUND,
            f"Source file was not found: {source_path}",
        )

    handler = get_local_extraction_handler(profile_name)
    if handler is None:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_UNSUPPORTED_PROFILE,
            f"Unsupported local extraction profile: {profile_name}",
        )

    if not handler.supports_file_type(file_type):
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_UNSUPPORTED_FILE_TYPE,
            f"{profile_name} cannot process {source_path.suffix or '(none)'} files",
        )

    result = handler.extract(request, source_path, started)
    validate_runtime_result(result)
    return result


def normalize_file_type(file_type: str | None) -> str:
    return (file_type or "").strip().lower().removeprefix(".")


def list_local_extraction_handlers() -> tuple[LocalExtractionHandler, ...]:
    return LOCAL_EXTRACTION_HANDLERS


def get_local_extraction_handler(profile_name: str | None) -> LocalExtractionHandler | None:
    normalized = (profile_name or "").strip()
    return LOCAL_EXTRACTION_HANDLER_BY_PROFILE.get(normalized)


def select_local_extraction_handler(
    file_type: str | None,
) -> LocalExtractionHandler | None:
    normalized = normalize_file_type(file_type)
    for handler in LOCAL_EXTRACTION_HANDLERS:
        if handler.supports_file_type(normalized):
            return handler
    return None


def select_local_extraction_profile_name(file_type: str | None) -> str | None:
    handler = select_local_extraction_handler(file_type)
    return handler.profile_name if handler else None


def _extract_markdown_source(
    request: ExtractionRuntimeRequest,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    source_text, failed_result = _read_utf8_text_source(source_path, started)
    if failed_result is not None:
        return failed_result
    return _extract_markdown(request, source_text, source_path, started)


def _extract_plain_text_source(
    request: ExtractionRuntimeRequest,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    source_text, failed_result = _read_utf8_text_source(source_path, started)
    if failed_result is not None:
        return failed_result
    return _extract_plain_text(request, source_text, source_path, started)


def _extract_pdf_text_source(
    request: ExtractionRuntimeRequest,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_PDF_READ_FAILED,
            f"pypdf is required for local PDF text extraction: {exc}",
        )

    try:
        reader = PdfReader(str(source_path))
        if reader.is_encrypted and not reader.decrypt(""):
            return _failed_result(
                started,
                ERROR_CODE_LOCAL_PDF_READ_FAILED,
                "Encrypted PDF could not be decrypted with an empty password",
            )
        page_count = len(reader.pages)
        page_texts = tuple(
            (page_no, _normalize_line_endings(page.extract_text() or "").strip())
            for page_no, page in enumerate(reader.pages, start=1)
        )
    except Exception as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_PDF_READ_FAILED,
            f"PDF text layer could not be read: {exc}",
        )

    extracted_page_texts = tuple(
        (page_no, page_text) for page_no, page_text in page_texts if page_text
    )
    if not extracted_page_texts:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_PDF_TEXT_LAYER_EMPTY,
            "PDF does not contain extractable text layer content",
        )

    artifact_text, blocks = _pdf_page_texts_to_artifact_and_blocks(extracted_page_texts)
    skipped_page_count = page_count - len(extracted_page_texts)
    library_version = _package_version(PDF_TEXT_LIBRARY_NAME)
    warnings = (
        (f"Skipped {skipped_page_count} PDF pages without extractable text",)
        if skipped_page_count
        else ()
    )
    metadata = {
        "source_path": str(source_path),
        "source_file_name": source_path.name,
        "parser_name": LOCAL_PDF_TEXT_EXTRACTOR_NAME,
        "parser_version": LOCAL_PDF_TEXT_EXTRACTOR_VERSION,
        "library": PDF_TEXT_LIBRARY_NAME,
        "library_version": library_version,
        "page_count": page_count,
        "extracted_page_count": len(extracted_page_texts),
        "block_count": len(blocks),
        "text_layer_only": True,
        "ocr_enabled": False,
        "options": request.options,
    }
    artifact = ExtractionRuntimeArtifact(
        artifact_type="normalized_markdown",
        content_text=artifact_text,
        content_hash=_hash_text(artifact_text),
        size_bytes=len(artifact_text.encode("utf-8")),
        metadata=metadata,
    )
    return ExtractionRuntimeResult(
        status="succeeded",
        artifacts=(artifact,),
        blocks=blocks,
        warnings=warnings,
        elapsed_ms=_elapsed_ms(started),
        runtime_metadata={
            "profile_name": LOCAL_PDF_TEXT_PROFILE_NAME,
            "extractor_name": LOCAL_PDF_TEXT_EXTRACTOR_NAME,
            "extractor_version": LOCAL_PDF_TEXT_EXTRACTOR_VERSION,
            "library": PDF_TEXT_LIBRARY_NAME,
            "library_version": library_version,
            "source_path": str(source_path),
            "page_count": page_count,
            "extracted_page_count": len(extracted_page_texts),
            "block_count": len(blocks),
            "text_layer_only": True,
            "ocr_enabled": False,
        },
    )


def _extract_docx_source(
    request: ExtractionRuntimeRequest,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    try:
        from docx import Document
    except ImportError as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_DOCX_READ_FAILED,
            f"python-docx is required for local DOCX extraction: {exc}",
        )

    try:
        document = Document(str(source_path))
        artifact_text, blocks, counts = _docx_document_to_artifact_and_blocks(document)
    except Exception as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_DOCX_READ_FAILED,
            f"DOCX content could not be read: {exc}",
        )

    if not blocks:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_DOCX_EMPTY,
            "DOCX does not contain extractable paragraph or table text",
        )

    library_version = _package_version(DOCX_LIBRARY_NAME)
    metadata = {
        "source_path": str(source_path),
        "source_file_name": source_path.name,
        "parser_name": LOCAL_DOCX_EXTRACTOR_NAME,
        "parser_version": LOCAL_DOCX_EXTRACTOR_VERSION,
        "library": DOCX_LIBRARY_NAME,
        "library_version": library_version,
        "paragraph_count": counts["paragraph_count"],
        "heading_count": counts["heading_count"],
        "table_count": counts["table_count"],
        "block_count": len(blocks),
        "preserve_headings": True,
        "preserve_tables": True,
        "options": request.options,
    }
    artifact = ExtractionRuntimeArtifact(
        artifact_type="normalized_markdown",
        content_text=artifact_text,
        content_hash=_hash_text(artifact_text),
        size_bytes=len(artifact_text.encode("utf-8")),
        metadata=metadata,
    )
    return ExtractionRuntimeResult(
        status="succeeded",
        artifacts=(artifact,),
        blocks=blocks,
        elapsed_ms=_elapsed_ms(started),
        runtime_metadata={
            "profile_name": LOCAL_DOCX_PROFILE_NAME,
            "extractor_name": LOCAL_DOCX_EXTRACTOR_NAME,
            "extractor_version": LOCAL_DOCX_EXTRACTOR_VERSION,
            "library": DOCX_LIBRARY_NAME,
            "library_version": library_version,
            "source_path": str(source_path),
            "paragraph_count": counts["paragraph_count"],
            "heading_count": counts["heading_count"],
            "table_count": counts["table_count"],
            "block_count": len(blocks),
            "preserve_headings": True,
            "preserve_tables": True,
        },
    )


def _extract_pptx_source(
    request: ExtractionRuntimeRequest,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_PPTX_READ_FAILED,
            f"python-pptx is required for local PPTX extraction: {exc}",
        )

    try:
        presentation = Presentation(str(source_path))
        artifact_text, blocks, counts = _pptx_presentation_to_artifact_and_blocks(presentation)
    except Exception as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_PPTX_READ_FAILED,
            f"PPTX content could not be read: {exc}",
        )

    if not blocks:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_PPTX_EMPTY,
            "PPTX does not contain extractable slide text or table text",
        )

    library_version = _package_version(PPTX_LIBRARY_NAME)
    metadata = {
        "source_path": str(source_path),
        "source_file_name": source_path.name,
        "parser_name": LOCAL_PPTX_EXTRACTOR_NAME,
        "parser_version": LOCAL_PPTX_EXTRACTOR_VERSION,
        "library": PPTX_LIBRARY_NAME,
        "library_version": library_version,
        "slide_count": counts["slide_count"],
        "text_shape_count": counts["text_shape_count"],
        "table_count": counts["table_count"],
        "block_count": len(blocks),
        "preserve_slide_boundaries": True,
        "preserve_tables": True,
        "options": request.options,
    }
    artifact = ExtractionRuntimeArtifact(
        artifact_type="normalized_markdown",
        content_text=artifact_text,
        content_hash=_hash_text(artifact_text),
        size_bytes=len(artifact_text.encode("utf-8")),
        metadata=metadata,
    )
    return ExtractionRuntimeResult(
        status="succeeded",
        artifacts=(artifact,),
        blocks=blocks,
        elapsed_ms=_elapsed_ms(started),
        runtime_metadata={
            "profile_name": LOCAL_PPTX_PROFILE_NAME,
            "extractor_name": LOCAL_PPTX_EXTRACTOR_NAME,
            "extractor_version": LOCAL_PPTX_EXTRACTOR_VERSION,
            "library": PPTX_LIBRARY_NAME,
            "library_version": library_version,
            "source_path": str(source_path),
            "slide_count": counts["slide_count"],
            "text_shape_count": counts["text_shape_count"],
            "table_count": counts["table_count"],
            "block_count": len(blocks),
            "preserve_slide_boundaries": True,
            "preserve_tables": True,
        },
    )


def _extract_xlsx_source(
    request: ExtractionRuntimeRequest,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_XLSX_READ_FAILED,
            f"openpyxl is required for local XLSX extraction: {exc}",
        )

    try:
        workbook = load_workbook(
            filename=str(source_path),
            read_only=True,
            data_only=True,
        )
        artifact_text, blocks, counts = _xlsx_workbook_to_artifact_and_blocks(workbook)
        workbook.close()
    except Exception as exc:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_XLSX_READ_FAILED,
            f"XLSX content could not be read: {exc}",
        )

    if not blocks:
        return _failed_result(
            started,
            ERROR_CODE_LOCAL_XLSX_EMPTY,
            "XLSX does not contain extractable worksheet cell values",
        )

    library_version = _package_version(XLSX_LIBRARY_NAME)
    metadata = {
        "source_path": str(source_path),
        "source_file_name": source_path.name,
        "parser_name": LOCAL_XLSX_EXTRACTOR_NAME,
        "parser_version": LOCAL_XLSX_EXTRACTOR_VERSION,
        "library": XLSX_LIBRARY_NAME,
        "library_version": library_version,
        "sheet_count": counts["sheet_count"],
        "extracted_sheet_count": counts["extracted_sheet_count"],
        "table_count": counts["table_count"],
        "cell_count": counts["cell_count"],
        "block_count": len(blocks),
        "preserve_sheet_boundaries": True,
        "emit_markdown_tables": True,
        "formulas_resolved_from_cached_values": True,
        "options": request.options,
    }
    artifact = ExtractionRuntimeArtifact(
        artifact_type="normalized_markdown",
        content_text=artifact_text,
        content_hash=_hash_text(artifact_text),
        size_bytes=len(artifact_text.encode("utf-8")),
        metadata=metadata,
    )
    return ExtractionRuntimeResult(
        status="succeeded",
        artifacts=(artifact,),
        blocks=blocks,
        elapsed_ms=_elapsed_ms(started),
        runtime_metadata={
            "profile_name": LOCAL_XLSX_PROFILE_NAME,
            "extractor_name": LOCAL_XLSX_EXTRACTOR_NAME,
            "extractor_version": LOCAL_XLSX_EXTRACTOR_VERSION,
            "library": XLSX_LIBRARY_NAME,
            "library_version": library_version,
            "source_path": str(source_path),
            "sheet_count": counts["sheet_count"],
            "extracted_sheet_count": counts["extracted_sheet_count"],
            "table_count": counts["table_count"],
            "cell_count": counts["cell_count"],
            "block_count": len(blocks),
            "preserve_sheet_boundaries": True,
            "emit_markdown_tables": True,
        },
    )


def _read_utf8_text_source(
    source_path: Path,
    started: float,
) -> tuple[str, ExtractionRuntimeResult | None]:
    try:
        source_text = _normalize_line_endings(source_path.read_text(encoding="utf-8"))
    except UnicodeDecodeError as exc:
        return "", _failed_result(
            started,
            ERROR_CODE_LOCAL_TEXT_DECODE_FAILED,
            f"Source file could not be decoded as UTF-8: {exc}",
        )

    if not source_text.strip():
        return "", _failed_result(
            started,
            ERROR_CODE_LOCAL_SOURCE_EMPTY,
            "Source file does not contain extractable text",
        )

    return source_text, None


def _pdf_page_texts_to_artifact_and_blocks(
    page_texts: tuple[tuple[int, str], ...],
) -> tuple[str, tuple[ExtractionRuntimeBlock, ...]]:
    parts: list[str] = []
    blocks: list[ExtractionRuntimeBlock] = []
    cursor = 0

    for page_no, page_text in page_texts:
        if parts:
            parts.append("\n\n")
            cursor += 2
        page_marker = f"<!-- page: {page_no} -->"
        parts.append(page_marker)
        cursor += len(page_marker)

        for paragraph_index, paragraph in enumerate(_split_pdf_paragraphs(page_text), start=1):
            parts.append("\n\n")
            cursor += 2
            char_start = cursor
            parts.append(paragraph)
            cursor += len(paragraph)
            blocks.append(
                ExtractionRuntimeBlock(
                    block_seq=len(blocks),
                    block_type="paragraph",
                    content_text=paragraph,
                    content_markdown=paragraph,
                    heading_path=(f"Page {page_no}",),
                    source_anchor={
                        "page_no": page_no,
                        "paragraph_index": paragraph_index,
                    },
                    page_no=page_no,
                    char_start=char_start,
                    char_end=cursor,
                    token_count=count_chunk_tokens(paragraph),
                    metadata={
                        "source": "pdf_text_layer",
                        "library": PDF_TEXT_LIBRARY_NAME,
                    },
                )
            )

    return "".join(parts), tuple(blocks)


def _split_pdf_paragraphs(page_text: str) -> tuple[str, ...]:
    paragraphs: list[str] = []
    current_lines: list[str] = []

    for line in page_text.splitlines():
        stripped = line.strip()
        if stripped:
            current_lines.append(stripped)
        elif current_lines:
            paragraphs.append(" ".join(current_lines))
            current_lines = []

    if current_lines:
        paragraphs.append(" ".join(current_lines))
    return tuple(paragraphs)


def _docx_document_to_artifact_and_blocks(document: Any) -> tuple[
    str,
    tuple[ExtractionRuntimeBlock, ...],
    dict[str, int],
]:
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    parts: list[str] = []
    blocks: list[ExtractionRuntimeBlock] = []
    latest_heading_seq_by_level: dict[int, int] = {}
    latest_heading_text_by_level: dict[int, str] = {}
    counts = {
        "paragraph_count": 0,
        "heading_count": 0,
        "table_count": 0,
    }
    cursor = 0
    paragraph_index = 0
    table_index = 0

    def append_block(
        *,
        block_type: str,
        content_text: str,
        content_markdown: str,
        heading_path: tuple[str, ...],
        source_anchor: dict[str, Any],
        parent_block_seq: int | None,
        metadata: dict[str, Any],
    ) -> None:
        nonlocal cursor
        if parts:
            parts.append("\n\n")
            cursor += 2
        char_start = cursor
        parts.append(content_markdown)
        cursor += len(content_markdown)
        blocks.append(
            ExtractionRuntimeBlock(
                block_seq=len(blocks),
                block_type=block_type,
                parent_block_seq=parent_block_seq,
                content_text=content_text,
                content_markdown=content_markdown,
                heading_path=heading_path,
                source_anchor=source_anchor,
                char_start=char_start,
                char_end=cursor,
                token_count=count_chunk_tokens(content_markdown),
                metadata=metadata,
            )
        )

    for body_index, child in enumerate(document.element.body.iterchildren()):
        if child.tag.endswith("}p"):
            paragraph = Paragraph(child, document)
            paragraph_text = _normalize_docx_text(paragraph.text)
            if not paragraph_text:
                continue
            paragraph_index += 1
            style_name = paragraph.style.name if paragraph.style is not None else None
            heading_level = _docx_heading_level(style_name)
            if heading_level is not None:
                parent_block_seq = latest_heading_seq_by_level.get(heading_level - 1)
                latest_heading_seq_by_level = {
                    level: seq
                    for level, seq in latest_heading_seq_by_level.items()
                    if level < heading_level
                }
                latest_heading_text_by_level = {
                    level: text
                    for level, text in latest_heading_text_by_level.items()
                    if level < heading_level
                }
                latest_heading_seq_by_level[heading_level] = len(blocks)
                latest_heading_text_by_level[heading_level] = paragraph_text
                heading_path = _docx_heading_path(latest_heading_text_by_level)
                block_type = "heading"
                content_markdown = f"{'#' * heading_level} {paragraph_text}"
                counts["heading_count"] += 1
            else:
                parent_block_seq = (
                    latest_heading_seq_by_level[max(latest_heading_seq_by_level)]
                    if latest_heading_seq_by_level
                    else None
                )
                heading_path = _docx_heading_path(latest_heading_text_by_level)
                block_type = "paragraph"
                content_markdown = paragraph_text
                counts["paragraph_count"] += 1

            append_block(
                block_type=block_type,
                content_text=paragraph_text,
                content_markdown=content_markdown,
                heading_path=heading_path,
                source_anchor={
                    "body_index": body_index,
                    "paragraph_index": paragraph_index,
                    "style_name": style_name,
                },
                parent_block_seq=parent_block_seq,
                metadata={
                    "source": "docx",
                    "style_name": style_name,
                    "heading_level": heading_level,
                },
            )
        elif child.tag.endswith("}tbl"):
            table = Table(child, document)
            rows = _docx_table_rows(table)
            if not rows:
                continue
            table_index += 1
            content_text = "\n".join("\t".join(row) for row in rows)
            content_markdown = _format_markdown_table(rows)
            parent_block_seq = (
                latest_heading_seq_by_level[max(latest_heading_seq_by_level)]
                if latest_heading_seq_by_level
                else None
            )
            append_block(
                block_type="table",
                content_text=content_text,
                content_markdown=content_markdown,
                heading_path=_docx_heading_path(latest_heading_text_by_level),
                source_anchor={
                    "body_index": body_index,
                    "table_index": table_index,
                },
                parent_block_seq=parent_block_seq,
                metadata={
                    "source": "docx",
                    "row_count": len(rows),
                    "column_count": max(len(row) for row in rows),
                },
            )
            counts["table_count"] += 1

    return "".join(parts), tuple(blocks), counts


def _pptx_presentation_to_artifact_and_blocks(presentation: Any) -> tuple[
    str,
    tuple[ExtractionRuntimeBlock, ...],
    dict[str, int],
]:
    parts: list[str] = []
    blocks: list[ExtractionRuntimeBlock] = []
    counts = {
        "slide_count": len(presentation.slides),
        "text_shape_count": 0,
        "table_count": 0,
    }
    cursor = 0
    marked_slide_numbers: set[int] = set()
    text_shape_index = 0
    table_index = 0

    def ensure_slide_marker(slide_no: int) -> None:
        nonlocal cursor
        if slide_no in marked_slide_numbers:
            return
        if parts:
            parts.append("\n\n")
            cursor += 2
        slide_marker = f"<!-- slide: {slide_no} -->"
        parts.append(slide_marker)
        cursor += len(slide_marker)
        marked_slide_numbers.add(slide_no)

    def append_block(
        *,
        slide_no: int,
        block_type: str,
        content_text: str,
        content_markdown: str,
        heading_path: tuple[str, ...],
        source_anchor: dict[str, Any],
        parent_block_seq: int | None,
        metadata: dict[str, Any],
    ) -> None:
        nonlocal cursor
        ensure_slide_marker(slide_no)
        parts.append("\n\n")
        cursor += 2
        char_start = cursor
        parts.append(content_markdown)
        cursor += len(content_markdown)
        blocks.append(
            ExtractionRuntimeBlock(
                block_seq=len(blocks),
                block_type=block_type,
                parent_block_seq=parent_block_seq,
                content_text=content_text,
                content_markdown=content_markdown,
                heading_path=heading_path,
                source_anchor=source_anchor,
                slide_no=slide_no,
                char_start=char_start,
                char_end=cursor,
                token_count=count_chunk_tokens(content_markdown),
                metadata=metadata,
            )
        )

    for slide_no, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_shape_id = title_shape.shape_id if title_shape is not None else None
        slide_title = _pptx_shape_text(title_shape) if title_shape is not None else ""
        slide_heading_path = (slide_title,) if slide_title else (f"Slide {slide_no}",)
        slide_heading_seq: int | None = None
        slide_text_shape_index = 0
        slide_table_index = 0

        for shape_index, shape in enumerate(slide.shapes, start=1):
            shape_name = getattr(shape, "name", None)
            if getattr(shape, "has_table", False):
                rows = _pptx_table_rows(shape.table)
                if not rows:
                    continue
                table_index += 1
                slide_table_index += 1
                counts["table_count"] += 1
                content_text = "\n".join("\t".join(row) for row in rows)
                append_block(
                    slide_no=slide_no,
                    block_type="table",
                    content_text=content_text,
                    content_markdown=_format_markdown_table(rows),
                    heading_path=slide_heading_path,
                    source_anchor={
                        "slide_no": slide_no,
                        "shape_index": shape_index,
                        "table_index": table_index,
                        "slide_table_index": slide_table_index,
                        "shape_name": shape_name,
                    },
                    parent_block_seq=slide_heading_seq,
                    metadata={
                        "source": "pptx",
                        "row_count": len(rows),
                        "column_count": max(len(row) for row in rows),
                        "shape_name": shape_name,
                    },
                )
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            content_text = _pptx_shape_text(shape)
            if not content_text:
                continue

            text_shape_index += 1
            slide_text_shape_index += 1
            counts["text_shape_count"] += 1
            is_title = shape.shape_id == title_shape_id
            if is_title:
                heading_text = _normalize_docx_text(content_text)
                block_type = "heading"
                content_text = heading_text
                content_markdown = f"# {heading_text}"
                parent_block_seq = None
                slide_heading_path = (heading_text,)
            else:
                block_type = "paragraph"
                content_markdown = content_text
                parent_block_seq = slide_heading_seq

            append_block(
                slide_no=slide_no,
                block_type=block_type,
                content_text=content_text,
                content_markdown=content_markdown,
                heading_path=slide_heading_path,
                source_anchor={
                    "slide_no": slide_no,
                    "shape_index": shape_index,
                    "text_shape_index": text_shape_index,
                    "slide_text_shape_index": slide_text_shape_index,
                    "shape_name": shape_name,
                },
                parent_block_seq=parent_block_seq,
                metadata={
                    "source": "pptx",
                    "shape_name": shape_name,
                    "is_title": is_title,
                },
            )
            if is_title:
                slide_heading_seq = blocks[-1].block_seq

    return "".join(parts), tuple(blocks), counts


def _pptx_shape_text(shape: Any) -> str:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return ""
    paragraphs = [
        _normalize_docx_text(paragraph.text)
        for paragraph in shape.text_frame.paragraphs
        if _normalize_docx_text(paragraph.text)
    ]
    return "\n".join(paragraphs)


def _pptx_table_rows(table: Any) -> tuple[tuple[str, ...], ...]:
    rows: list[tuple[str, ...]] = []
    for table_row in table.rows:
        values = tuple(_normalize_docx_text(cell.text) for cell in table_row.cells)
        if any(values):
            rows.append(values)
    return tuple(rows)


def _xlsx_workbook_to_artifact_and_blocks(workbook: Any) -> tuple[
    str,
    tuple[ExtractionRuntimeBlock, ...],
    dict[str, int],
]:
    from openpyxl.utils import get_column_letter

    parts: list[str] = []
    blocks: list[ExtractionRuntimeBlock] = []
    counts = {
        "sheet_count": len(workbook.worksheets),
        "extracted_sheet_count": 0,
        "table_count": 0,
        "cell_count": 0,
    }
    cursor = 0

    def append_block(
        *,
        block_type: str,
        content_text: str,
        content_markdown: str,
        heading_path: tuple[str, ...],
        source_anchor: dict[str, Any],
        parent_block_seq: int | None,
        metadata: dict[str, Any],
        sheet_name: str,
        cell_range: str | None = None,
    ) -> None:
        nonlocal cursor
        if parts:
            parts.append("\n\n")
            cursor += 2
        char_start = cursor
        parts.append(content_markdown)
        cursor += len(content_markdown)
        blocks.append(
            ExtractionRuntimeBlock(
                block_seq=len(blocks),
                block_type=block_type,
                parent_block_seq=parent_block_seq,
                content_text=content_text,
                content_markdown=content_markdown,
                heading_path=heading_path,
                source_anchor=source_anchor,
                sheet_name=sheet_name,
                cell_range=cell_range,
                char_start=char_start,
                char_end=cursor,
                token_count=count_chunk_tokens(content_markdown),
                metadata=metadata,
            )
        )

    for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
        rows, bounds, cell_count = _xlsx_non_empty_rows(worksheet)
        if not rows or bounds is None:
            continue

        counts["extracted_sheet_count"] += 1
        counts["table_count"] += 1
        counts["cell_count"] += cell_count
        min_row, min_column, max_row, max_column = bounds
        cell_range = (
            f"{get_column_letter(min_column)}{min_row}:" f"{get_column_letter(max_column)}{max_row}"
        )
        sheet_marker = f"<!-- sheet: {worksheet.title} -->"
        if parts:
            parts.append("\n\n")
            cursor += 2
        parts.append(sheet_marker)
        cursor += len(sheet_marker)

        heading_text = worksheet.title
        heading_markdown = f"# {heading_text}"
        append_block(
            block_type="heading",
            content_text=heading_text,
            content_markdown=heading_markdown,
            heading_path=(heading_text,),
            source_anchor={
                "sheet_index": sheet_index,
                "sheet_name": worksheet.title,
            },
            parent_block_seq=None,
            metadata={
                "source": "xlsx",
                "sheet_index": sheet_index,
            },
            sheet_name=worksheet.title,
        )
        append_block(
            block_type="table",
            content_text="\n".join("\t".join(row) for row in rows),
            content_markdown=_format_markdown_table(rows),
            heading_path=(heading_text,),
            source_anchor={
                "sheet_index": sheet_index,
                "sheet_name": worksheet.title,
                "cell_range": cell_range,
                "table_index": counts["table_count"],
            },
            parent_block_seq=blocks[-1].block_seq,
            metadata={
                "source": "xlsx",
                "row_count": len(rows),
                "column_count": max(len(row) for row in rows),
                "cell_count": cell_count,
            },
            sheet_name=worksheet.title,
            cell_range=cell_range,
        )

    return "".join(parts), tuple(blocks), counts


def _xlsx_non_empty_rows(worksheet: Any) -> tuple[
    tuple[tuple[str, ...], ...],
    tuple[int, int, int, int] | None,
    int,
]:
    rows: list[tuple[int, tuple[str, ...]]] = []
    min_row: int | None = None
    max_row: int | None = None
    min_column: int | None = None
    max_column: int | None = None
    cell_count = 0

    for row_index, row in enumerate(worksheet.iter_rows(), start=1):
        values = tuple(_normalize_xlsx_cell_value(cell.value) for cell in row)
        non_empty_positions = [
            column_index for column_index, value in enumerate(values, start=1) if value
        ]
        if not non_empty_positions:
            continue

        rows.append((row_index, values))
        cell_count += len(non_empty_positions)
        row_min_column = min(non_empty_positions)
        row_max_column = max(non_empty_positions)
        min_row = row_index if min_row is None else min(min_row, row_index)
        max_row = row_index if max_row is None else max(max_row, row_index)
        min_column = row_min_column if min_column is None else min(min_column, row_min_column)
        max_column = row_max_column if max_column is None else max(max_column, row_max_column)

    if min_row is None or min_column is None or max_row is None or max_column is None:
        return (), None, 0

    normalized_rows = tuple(
        tuple(values[column_index - 1] for column_index in range(min_column, max_column + 1))
        for _, values in rows
    )
    return normalized_rows, (min_row, min_column, max_row, max_column), cell_count


def _normalize_xlsx_cell_value(value: Any) -> str:
    if value is None:
        return ""
    return _normalize_docx_text(str(value))


def _normalize_docx_text(text: str) -> str:
    return " ".join(text.replace("\r", "\n").split())


def _docx_heading_level(style_name: str | None) -> int | None:
    normalized = (style_name or "").strip().lower()
    if not normalized.startswith("heading "):
        return None
    raw_level = normalized.removeprefix("heading ").strip()
    if not raw_level.isdigit():
        return None
    level = int(raw_level)
    return level if level > 0 else None


def _docx_heading_path(latest_heading_text_by_level: dict[int, str]) -> tuple[str, ...]:
    return tuple(
        latest_heading_text_by_level[level] for level in sorted(latest_heading_text_by_level)
    )


def _docx_table_rows(table: Any) -> tuple[tuple[str, ...], ...]:
    rows: list[tuple[str, ...]] = []
    for table_row in table.rows:
        values = tuple(_normalize_docx_text(cell.text) for cell in table_row.cells)
        if any(values):
            rows.append(values)
    return tuple(rows)


def _format_markdown_table(rows: tuple[tuple[str, ...], ...]) -> str:
    column_count = max(len(row) for row in rows)
    normalized_rows = [
        tuple(row[index] if index < len(row) else "" for index in range(column_count))
        for row in rows
    ]
    header = normalized_rows[0]
    separator = tuple("---" for _ in range(column_count))
    body_rows = normalized_rows[1:]
    return "\n".join(
        [_markdown_table_row(header), _markdown_table_row(separator)]
        + [_markdown_table_row(row) for row in body_rows]
    )


def _markdown_table_row(row: tuple[str, ...]) -> str:
    return "| " + " | ".join(value.replace("|", "\\|") for value in row) + " |"


LOCAL_EXTRACTION_HANDLERS = (
    LocalExtractionHandler(
        profile_name=LOCAL_MARKDOWN_PROFILE_NAME,
        parser_name=PARSER_NAME_MARKDOWN,
        parser_version=PARSER_VERSION_MARKDOWN,
        extractor_name="local_markdown",
        extractor_version=PARSER_VERSION_MARKDOWN,
        supported_file_types=("md",),
        extract=_extract_markdown_source,
    ),
    LocalExtractionHandler(
        profile_name=LOCAL_PLAIN_TEXT_PROFILE_NAME,
        parser_name=LOCAL_PLAIN_TEXT_EXTRACTOR_NAME,
        parser_version=LOCAL_PLAIN_TEXT_EXTRACTOR_VERSION,
        extractor_name=LOCAL_PLAIN_TEXT_EXTRACTOR_NAME,
        extractor_version=LOCAL_PLAIN_TEXT_EXTRACTOR_VERSION,
        supported_file_types=("txt", "text"),
        extract=_extract_plain_text_source,
    ),
    LocalExtractionHandler(
        profile_name=LOCAL_PDF_TEXT_PROFILE_NAME,
        parser_name=LOCAL_PDF_TEXT_EXTRACTOR_NAME,
        parser_version=LOCAL_PDF_TEXT_EXTRACTOR_VERSION,
        extractor_name=LOCAL_PDF_TEXT_EXTRACTOR_NAME,
        extractor_version=LOCAL_PDF_TEXT_EXTRACTOR_VERSION,
        supported_file_types=("pdf",),
        extract=_extract_pdf_text_source,
    ),
    LocalExtractionHandler(
        profile_name=LOCAL_DOCX_PROFILE_NAME,
        parser_name=LOCAL_DOCX_EXTRACTOR_NAME,
        parser_version=LOCAL_DOCX_EXTRACTOR_VERSION,
        extractor_name=LOCAL_DOCX_EXTRACTOR_NAME,
        extractor_version=LOCAL_DOCX_EXTRACTOR_VERSION,
        supported_file_types=("docx",),
        extract=_extract_docx_source,
    ),
    LocalExtractionHandler(
        profile_name=LOCAL_PPTX_PROFILE_NAME,
        parser_name=LOCAL_PPTX_EXTRACTOR_NAME,
        parser_version=LOCAL_PPTX_EXTRACTOR_VERSION,
        extractor_name=LOCAL_PPTX_EXTRACTOR_NAME,
        extractor_version=LOCAL_PPTX_EXTRACTOR_VERSION,
        supported_file_types=("pptx",),
        extract=_extract_pptx_source,
    ),
    LocalExtractionHandler(
        profile_name=LOCAL_XLSX_PROFILE_NAME,
        parser_name=LOCAL_XLSX_EXTRACTOR_NAME,
        parser_version=LOCAL_XLSX_EXTRACTOR_VERSION,
        extractor_name=LOCAL_XLSX_EXTRACTOR_NAME,
        extractor_version=LOCAL_XLSX_EXTRACTOR_VERSION,
        supported_file_types=("xlsx",),
        extract=_extract_xlsx_source,
    ),
)

LOCAL_EXTRACTION_HANDLER_BY_PROFILE = {
    handler.profile_name: handler for handler in LOCAL_EXTRACTION_HANDLERS
}

SUPPORTED_LOCAL_PROFILE_SUFFIXES = {
    handler.profile_name: {f".{file_type}" for file_type in handler.supported_file_types}
    for handler in LOCAL_EXTRACTION_HANDLERS
}


def persist_extraction_runtime_result(
    database_url: str,
    request: ExtractionRuntimeRequest,
    result: ExtractionRuntimeResult,
) -> PersistedExtractionRuntimeResult:
    validate_runtime_request(request)
    validate_runtime_result(result)
    with connect(database_url) as connection:
        return persist_extraction_runtime_result_in_connection(connection, request, result)


def persist_extraction_runtime_result_in_connection(
    connection: Connection,
    request: ExtractionRuntimeRequest,
    result: ExtractionRuntimeResult,
) -> PersistedExtractionRuntimeResult:
    validate_runtime_request(request)
    validate_runtime_result(result)
    profile = get_extraction_profile_in_connection(
        connection,
        request.extraction_profile_name,
    )
    if profile is None:
        raise InvalidIngestionArtifactError(
            f"Extraction profile was not found: {request.extraction_profile_name}"
        )
    if result.blocks and request.document_id is None:
        raise InvalidIngestionArtifactError("document_id is required when persisting blocks")

    run = create_extraction_run_in_connection(
        connection,
        ExtractionRunInput(
            file_id=request.file_id,
            document_id=request.document_id,
            extraction_profile_name=profile.extraction_profile_name,
            status=result.status,
            provider_mode=profile.provider_mode,
            extractor_name=profile.extractor_name,
            extractor_version=profile.extractor_version,
            elapsed_ms=result.elapsed_ms,
            warning_count=len(result.warnings),
            error_count=len(result.errors),
            error_code=(
                str(result.runtime_metadata.get("error_code"))
                if result.runtime_metadata.get("error_code")
                else None
            ),
            error_message=result.errors[0] if result.errors else None,
            runtime_metadata=result.runtime_metadata,
        ),
    )

    artifacts = tuple(
        create_extraction_artifact_in_connection(
            connection,
            ExtractionArtifactInput(
                extraction_run_id=run.extraction_run_id,
                file_id=request.file_id,
                document_id=request.document_id,
                artifact_type=artifact.artifact_type,
                content_text=artifact.content_text,
                storage_path=artifact.storage_path,
                content_hash=artifact.content_hash,
                size_bytes=artifact.size_bytes,
                language=artifact.language,
                metadata=artifact.metadata,
            ),
        )
        for artifact in result.artifacts
    )
    if result.blocks and not artifacts:
        raise InvalidIngestionArtifactError("runtime blocks require a persisted artifact")

    block_records = _persist_blocks(
        connection,
        request.document_id,
        artifacts[0],
        result.blocks,
    )
    return PersistedExtractionRuntimeResult(
        run=run,
        artifacts=artifacts,
        blocks=block_records,
    )


def _extract_markdown(
    request: ExtractionRuntimeRequest,
    source_text: str,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    parser = MarkdownParser()
    parsed_document = parser.parse_text(source_text, source_path=str(source_path))
    line_offsets = _line_offsets(source_text)
    blocks = _markdown_blocks_to_runtime_blocks(parsed_document.blocks, line_offsets)
    artifact = ExtractionRuntimeArtifact(
        artifact_type="normalized_markdown",
        content_text=source_text,
        content_hash=_hash_text(source_text),
        size_bytes=len(source_text.encode("utf-8")),
        metadata={
            "source_path": str(source_path),
            "source_file_name": source_path.name,
            "parser_name": parsed_document.parser_name,
            "parser_version": parsed_document.parser_version,
            "line_count": parsed_document.line_count,
            "block_count": len(blocks),
            "options": request.options,
        },
    )
    return ExtractionRuntimeResult(
        status="succeeded",
        artifacts=(artifact,),
        blocks=blocks,
        elapsed_ms=_elapsed_ms(started),
        runtime_metadata={
            "profile_name": LOCAL_MARKDOWN_PROFILE_NAME,
            "extractor_name": PARSER_NAME_MARKDOWN,
            "extractor_version": PARSER_VERSION_MARKDOWN,
            "source_path": str(source_path),
            "line_count": parsed_document.line_count,
            "block_count": len(blocks),
        },
    )


def _extract_plain_text(
    request: ExtractionRuntimeRequest,
    source_text: str,
    source_path: Path,
    started: float,
) -> ExtractionRuntimeResult:
    blocks = tuple(_plain_text_blocks(source_text))
    artifact = ExtractionRuntimeArtifact(
        artifact_type="plain_text",
        content_text=source_text,
        content_hash=_hash_text(source_text),
        size_bytes=len(source_text.encode("utf-8")),
        metadata={
            "source_path": str(source_path),
            "source_file_name": source_path.name,
            "line_count": len(source_text.splitlines()),
            "block_count": len(blocks),
            "options": request.options,
        },
    )
    return ExtractionRuntimeResult(
        status="succeeded",
        artifacts=(artifact,),
        blocks=blocks,
        elapsed_ms=_elapsed_ms(started),
        runtime_metadata={
            "profile_name": LOCAL_PLAIN_TEXT_PROFILE_NAME,
            "extractor_name": LOCAL_PLAIN_TEXT_EXTRACTOR_NAME,
            "extractor_version": LOCAL_PLAIN_TEXT_EXTRACTOR_VERSION,
            "source_path": str(source_path),
            "line_count": len(source_text.splitlines()),
            "block_count": len(blocks),
        },
    )


def _markdown_blocks_to_runtime_blocks(
    parsed_blocks: tuple[ParsedBlock, ...],
    line_offsets: tuple[int, ...],
) -> tuple[ExtractionRuntimeBlock, ...]:
    runtime_blocks: list[ExtractionRuntimeBlock] = []
    latest_heading_seq_by_level: dict[int, int] = {}

    for block_seq, block in enumerate(parsed_blocks):
        heading_level = int(block.metadata.get("level", len(block.heading_path) or 1))
        parent_block_seq: int | None = None
        if block.block_type == "heading":
            parent_block_seq = latest_heading_seq_by_level.get(heading_level - 1)
            latest_heading_seq_by_level = {
                level: seq
                for level, seq in latest_heading_seq_by_level.items()
                if level < heading_level
            }
            latest_heading_seq_by_level[heading_level] = block_seq
        elif latest_heading_seq_by_level:
            parent_block_seq = latest_heading_seq_by_level[max(latest_heading_seq_by_level)]

        runtime_blocks.append(
            _markdown_block_to_runtime_block(
                block,
                block_seq,
                parent_block_seq,
                line_offsets,
            )
        )
    return tuple(runtime_blocks)


def _markdown_block_to_runtime_block(
    block: ParsedBlock,
    block_seq: int,
    parent_block_seq: int | None,
    line_offsets: tuple[int, ...],
) -> ExtractionRuntimeBlock:
    block_type = "code" if block.block_type == "code_block" else block.block_type
    content_markdown = _format_markdown_block(block)
    char_start, char_end = _char_range_for_lines(line_offsets, block.start_line, block.end_line)
    source_anchor: dict[str, Any] = {
        "start_line": block.start_line,
        "end_line": block.end_line,
    }
    if block.block_type == "table":
        source_anchor["table_index"] = block_seq
    return ExtractionRuntimeBlock(
        block_seq=block_seq,
        block_type=block_type,
        parent_block_seq=parent_block_seq,
        content_text=block.text,
        content_markdown=content_markdown,
        heading_path=block.heading_path,
        source_anchor=source_anchor,
        char_start=char_start,
        char_end=char_end,
        token_count=count_chunk_tokens(content_markdown),
        metadata=block.metadata,
    )


def _plain_text_blocks(source_text: str) -> list[ExtractionRuntimeBlock]:
    lines = source_text.splitlines(keepends=True)
    blocks: list[ExtractionRuntimeBlock] = []
    paragraph_lines: list[str] = []
    paragraph_start_line: int | None = None
    paragraph_start_char: int | None = None
    char_offset = 0

    for index, raw_line in enumerate(lines, start=1):
        line_text = raw_line.rstrip("\n")
        if line_text.strip():
            if paragraph_start_line is None:
                paragraph_start_line = index
                paragraph_start_char = char_offset
            paragraph_lines.append(line_text.strip())
        elif paragraph_lines:
            blocks.append(
                _build_plain_text_block(
                    blocks,
                    paragraph_lines,
                    paragraph_start_line,
                    index - 1,
                    paragraph_start_char,
                    char_offset,
                )
            )
            paragraph_lines = []
            paragraph_start_line = None
            paragraph_start_char = None
        char_offset += len(raw_line)

    if paragraph_lines:
        blocks.append(
            _build_plain_text_block(
                blocks,
                paragraph_lines,
                paragraph_start_line,
                len(lines),
                paragraph_start_char,
                len(source_text),
            )
        )
    return blocks


def _build_plain_text_block(
    blocks: list[ExtractionRuntimeBlock],
    paragraph_lines: list[str],
    start_line: int | None,
    end_line: int,
    start_char: int | None,
    end_char: int,
) -> ExtractionRuntimeBlock:
    content_text = " ".join(paragraph_lines)
    return ExtractionRuntimeBlock(
        block_seq=len(blocks),
        block_type="paragraph",
        content_text=content_text,
        content_markdown=content_text,
        source_anchor={
            "start_line": start_line or 1,
            "end_line": end_line,
        },
        char_start=start_char if start_char is not None else 0,
        char_end=end_char,
        token_count=count_chunk_tokens(content_text),
    )


def _persist_blocks(
    connection: Connection,
    document_id: int | None,
    artifact: ExtractionArtifactRecord,
    runtime_blocks: tuple[ExtractionRuntimeBlock, ...],
) -> tuple[DocumentBlockRecord, ...]:
    if document_id is None:
        return ()
    block_id_by_seq: dict[int, int] = {}
    records: list[DocumentBlockRecord] = []
    for block in runtime_blocks:
        parent_block_id = (
            block_id_by_seq.get(block.parent_block_seq)
            if block.parent_block_seq is not None
            else None
        )
        record = create_document_block_in_connection(
            connection,
            DocumentBlockInput(
                artifact_id=artifact.artifact_id,
                document_id=document_id,
                parent_block_id=parent_block_id,
                block_seq=block.block_seq,
                block_type=block.block_type,
                content_text=block.content_text,
                content_markdown=block.content_markdown,
                heading_path=block.heading_path,
                source_anchor=block.source_anchor,
                page_no=block.page_no,
                slide_no=block.slide_no,
                sheet_name=block.sheet_name,
                cell_range=block.cell_range,
                char_start=block.char_start,
                char_end=block.char_end,
                token_count=block.token_count,
                metadata=block.metadata,
            ),
        )
        block_id_by_seq[block.block_seq] = record.block_id
        records.append(record)
    return tuple(records)


def _format_markdown_block(block: ParsedBlock) -> str:
    if block.block_type == "heading":
        level = int(block.metadata.get("level", len(block.heading_path) or 1))
        return f"{'#' * level} {block.text}"
    if block.block_type == "code_block":
        language = block.metadata.get("language")
        fence = f"```{language}" if language else "```"
        return f"{fence}\n{block.text}\n```"
    return block.text


def _line_offsets(text: str) -> tuple[int, ...]:
    offsets = [0]
    position = 0
    for line in text.splitlines(keepends=True):
        position += len(line)
        offsets.append(position)
    if not text.endswith("\n"):
        offsets[-1] = len(text)
    return tuple(offsets)


def _char_range_for_lines(
    line_offsets: tuple[int, ...],
    start_line: int,
    end_line: int,
) -> tuple[int, int]:
    start_index = max(start_line - 1, 0)
    end_index = min(end_line, len(line_offsets) - 1)
    return line_offsets[start_index], line_offsets[end_index]


def _failed_result(
    started: float,
    error_code: str,
    error_message: str,
) -> ExtractionRuntimeResult:
    result = ExtractionRuntimeResult(
        status="failed",
        errors=(error_message,),
        elapsed_ms=_elapsed_ms(started),
        runtime_metadata={"error_code": error_code},
    )
    validate_runtime_result(result)
    return result


def _normalize_line_endings(text: str) -> str:
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _hash_text(text: str) -> str:
    return sha256(text.encode("utf-8")).hexdigest()


def _package_version(package_name: str) -> str:
    try:
        return version(package_name)
    except PackageNotFoundError:
        return "unknown"


def _elapsed_ms(started: float) -> int:
    return max(int((perf_counter() - started) * 1000), 0)
